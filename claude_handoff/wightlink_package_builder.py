from __future__ import annotations

import csv
import json
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any, Mapping, Sequence
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation


ASSETS_DIR = Path(__file__).resolve().parent / "wightlink_assets"
DEFAULT_REFERENCE_PPTX_PATH = ASSETS_DIR / "reference_deck_exported_from_google_slides.pptx"
DEFAULT_REFERENCE_DECK_URL = "https://docs.google.com/presentation/d/1QBixtNRD_N_qZ6J56RyfM1vvrTudrZUowZIMdBoGhso"
REFERENCE_SLIDE_COUNT = 27
TARGET_OUTPUT_SLIDE_COUNT = 22
PLAN_SOURCE_URL = "https://docs.google.com/spreadsheets/d/18w3DWmDtJ5plGWuzHjGNnnQ-M-uvyCrhIBKeaSTOl9Q/edit?gid=596378878#gid=596378878"
MONTHLY_REFERENCE_PPTX_FILENAME = "wightlink_qbr_visual_reference_only.pptx"
EXPECTED_TREND_TERMS = (
    "Wightlink Ferries",
    "Isle of Wight Ferry",
    "Isle of Wight Holidays",
)

PERIOD_RE = re.compile(r"\bQ[1-4]\s+\d{4}\s*\([^)]+\)")
MONTHLY_PERIOD_RE = re.compile(r"\b[A-Z][a-z]{2}\s+\d{4}\s*\(YTD\s+Jan\s+-\s+[A-Z][a-z]{2}\s+\d{4}\)")
QUARTER_SHORT_RE = re.compile(r"\bQ[1-4]\s+\d{4}\b")
MONTH_SHORT_RE = re.compile(r"\b[A-Z][a-z]{2}\s+\d{4}\b")
SECTION_RE = re.compile(r"^\[(?P<title>[^\]]+)\]\s*$")
KPI_RE = re.compile(r"^-\s*(?P<label>Cost|Purchases|Purchase Revenue|CPA|ROAS|AOV):\s*(?P<value>.+?)(?:\s+\(.*\))?$")

EXPECTED_EXACT_SECTIONS = (
    "Cover",
    "Agenda",
    "Trends",
    "Google Trends YTD - Wightlink Ferries",
    "Google Trends YTD - Isle of Wight Ferry",
    "Google Trends YTD - Isle of Wight Holidays",
    "Auction Insights - Generic",
    "Auction Insights - Brand",
    "Auction Insights - Red Funnel Quarter",
    "Performance",
    "All Performance Summary",
    "All Performance Purchases YoY",
    "Brand Performance Summary",
    "Brand Performance Monthly Purchases and Revenue",
    "Brand Monthly Breakdown YTD",
    "Generics Performance Summary",
    "Generics Performance Monthly Purchases and Revenue",
    "Generics Monthly Breakdown YTD",
    "PMax Performance Summary",
    "PMax Performance Monthly Purchases and Revenue",
    "PMax Performance Summary YTD",
    "Closing",
)


@dataclass(frozen=True)
class WightlinkRawInput:
    source_name: str
    archive_name: str
    role: str
    required: bool
    payload: bytes
    trend_query: str | None = None
    trend_period: str | None = None


@dataclass(frozen=True)
class WightlinkSourceSection:
    index: int
    title: str
    line: int
    slide: int | None = None


def is_wightlink_qbr(client_id: str, report_mode: str = "quarterly") -> bool:
    return client_id == "wightlink" and report_mode == "quarterly"


def is_wightlink_report(client_id: str, report_mode: str = "quarterly") -> bool:
    return client_id == "wightlink" and report_mode in {"quarterly", "monthly"}


def build_wightlink_claude_handoff_package(
    *,
    report_text: str,
    prompt_text: str,
    generated_pptx: bytes | Path | str,
    performance_csv: bytes | Path | str,
    auction_csv: bytes | Path | str | None = None,
    trend_csv_files: Sequence[bytes | Path | str] | None = None,
    trend_ytd_current_csv_files: Sequence[bytes | Path | str] | None = None,
    trend_ytd_previous_csv_files: Sequence[bytes | Path | str] | None = None,
    red_funnel_auction_csv: bytes | Path | str | None = None,
    red_funnel_prior_auction_csv: bytes | Path | str | None = None,
    plan_book_csv: bytes | Path | str | None = None,
    reference_pptx: bytes | Path | str | None = DEFAULT_REFERENCE_PPTX_PATH,
    reference_deck_url: str = DEFAULT_REFERENCE_DECK_URL,
    source_generation_manifest: Mapping[str, Any] | bytes | Path | str | None = None,
    generated_at: datetime | None = None,
    report_mode: str = "quarterly",
) -> tuple[bytes, dict]:
    is_monthly = report_mode == "monthly"
    generated_pptx_bytes = _read_payload(generated_pptx)
    reference_pptx_bytes = _read_optional_payload(reference_pptx)
    source_generation_payload, source_generation_json = _read_optional_json_manifest(source_generation_manifest)
    raw_inputs = (
        build_monthly_raw_inputs(performance_csv=performance_csv, plan_book_csv=plan_book_csv)
        if is_monthly
        else build_raw_inputs(
            performance_csv=performance_csv,
            auction_csv=auction_csv,
            trend_csv_files=trend_csv_files or [],
            trend_ytd_current_csv_files=trend_ytd_current_csv_files or [],
            trend_ytd_previous_csv_files=trend_ytd_previous_csv_files or [],
            red_funnel_auction_csv=red_funnel_auction_csv,
            red_funnel_prior_auction_csv=red_funnel_prior_auction_csv,
            plan_book_csv=plan_book_csv,
        )
    )

    source_sections = parse_source_sections(report_text)
    period_label = extract_period_label(report_text)
    quarter_short = extract_period_short(period_label)
    ytd_period_label, previous_ytd_period_label = extract_ytd_period_labels(report_text, quarter_short)
    headline_kpis = extract_headline_kpis(report_text)
    trend_queries = [raw_input.trend_query for raw_input in raw_inputs if raw_input.trend_query]
    trend_sources = build_trend_sources_manifest(raw_inputs)
    plan_metrics = extract_plan_metrics(report_text)
    streamlit_slide_count = count_pptx_slides(generated_pptx_bytes) or _max_report_slide_number(source_sections)
    missing_sections = find_monthly_missing_sections(source_sections) if is_monthly else find_missing_sections(source_sections)
    extra_sections = find_monthly_extra_sections(source_sections) if is_monthly else find_extra_sections(source_sections)

    warnings: list[str] = []
    if missing_sections:
        warnings.append(f"Missing expected sections: {', '.join(missing_sections)}")
    if extra_sections:
        warnings.append(f"Extra source sections found: {', '.join(extra_sections)}")
    if _has_filename_based_trend_titles(source_sections):
        warnings.append("Trend section titles appear filename-derived; use Google Trends CSV header query names instead.")
    if reference_pptx_bytes is None:
        warnings.append(f"{_reference_pptx_filename(is_monthly)} is not included; Claude needs it for full visual fidelity.")
    for required_name in _missing_required_raw_inputs(raw_inputs, is_monthly=is_monthly):
        warnings.append(f"Missing raw input: {required_name}")
    if is_monthly and plan_book_csv is None:
        warnings.append("Missing Wightlink plan source; monthly KPI Plan lines may be unavailable.")
    if not is_monthly:
        warnings.extend(_build_v2_source_warnings(trend_sources, raw_inputs, plan_book_csv))

    variables = {
        "client_display_name": "Wightlink",
        "report_family": "Wightlink PPC Monthly" if is_monthly else "Wightlink PPC QBR",
        "report_mode": report_mode,
        "period_label": period_label,
        "quarter_short": quarter_short,
        "ytd_period_label": ytd_period_label,
        "previous_ytd_period_label": previous_ytd_period_label,
        "reference_deck_url": reference_deck_url,
        "reference_slide_count": str(REFERENCE_SLIDE_COUNT),
        "target_output_slide_count": str(streamlit_slide_count if is_monthly else TARGET_OUTPUT_SLIDE_COUNT),
        "streamlit_slide_count": str(streamlit_slide_count),
        "streamlit_pptx_filename": "wightlink_monthly_streamlit_output.pptx" if is_monthly else "wightlink_streamlit_output.pptx",
        "reference_pptx_filename": _reference_pptx_filename(is_monthly),
        "performance_original_filename": _source_name(performance_csv),
        "auction_original_filename": _source_name(auction_csv) if auction_csv is not None else "not supplied",
        "plan_original_filename": _source_name(plan_book_csv) if plan_book_csv is not None else "not supplied",
        "red_funnel_quarter_auction_original_filename": _source_name(red_funnel_auction_csv) if red_funnel_auction_csv is not None else "not supplied",
        "red_funnel_prior_auction_original_filename": _source_name(red_funnel_prior_auction_csv) if red_funnel_prior_auction_csv is not None else "not supplied",
        "headline_cost": headline_kpis["cost"],
        "headline_purchases": headline_kpis["purchases"],
        "headline_purchase_revenue": headline_kpis["purchase_revenue"],
        "headline_cpa": headline_kpis["cpa"],
        "headline_roas": headline_kpis["roas"],
        "headline_aov": headline_kpis["aov"],
        "trend_queries": ", ".join(trend_queries) if trend_queries else "none supplied",
        "has_plan_source": "true" if plan_book_csv else "false",
        "plan_source_url": PLAN_SOURCE_URL,
    }
    variables.update(_trend_template_variables(trend_sources))

    if is_monthly:
        readme_text = build_monthly_readme_text(variables)
        claude_prompt_text = build_monthly_claude_prompt_text(variables)
        slide_mapping_text = build_monthly_slide_mapping_text(source_sections)
    else:
        readme_text = render_asset_template("README_FOR_CLAUDE_TEMPLATE.txt", variables)
        claude_prompt_text = render_asset_template("CLAUDE_PROMPT_TEMPLATE.txt", variables)
        slide_mapping_text = render_asset_template("SLIDE_MAPPING_WIGHTLINK_QBR_TEMPLATE.csv", variables)
    source_index_text = build_source_section_index(
        report_text=report_text,
        source_sections=source_sections,
        period_label=period_label,
        missing_sections=missing_sections,
        extra_sections=extra_sections,
        trend_queries=trend_queries,
        ytd_period_label=ytd_period_label,
        previous_ytd_period_label=previous_ytd_period_label,
        trend_sources=trend_sources,
        report_mode=report_mode,
    )
    input_files_manifest_text = build_raw_inputs_manifest(
        variables=variables,
        raw_inputs=raw_inputs,
        missing_required_inputs=_missing_required_raw_inputs(raw_inputs, is_monthly=is_monthly),
    )
    if is_monthly:
        qa_checklist_text = build_monthly_qa_checklist_text(variables)
        chart_qa_text = build_monthly_chart_qa_addendum_text(source_sections, variables)
        reference_outline_text = build_monthly_reference_outline_text(variables)
    else:
        qa_checklist_text = render_asset_template("QA_CHECKLIST_TEMPLATE.txt", variables).rstrip() + "\n\n" + read_asset_text("QA_CHECKLIST_V2_ADDENDUM.txt")
        chart_qa_text = read_asset_text("CHART_QA_ADDENDUM_FOR_CLAUDE.txt")
        reference_outline_text = render_asset_template("REFERENCE_DECK_OUTLINE.txt", variables)

    if reference_pptx_bytes is None:
        reference_note = (
            "\n\nReference deck availability note\n"
            f"- {_reference_pptx_filename(is_monthly)} is not included in this package.\n"
            "- Export or attach the Wightlink visual reference before asking Claude to complete the deck with full visual fidelity.\n"
        )
        readme_text += reference_note
        claude_prompt_text += reference_note

    if source_generation_payload is not None:
        source_note = (
            "\n\nAPI source-generation note\n"
            "- Performance and/or Google Trends source files were generated by the Streamlit API Source Test page.\n"
            "- See SOURCE_GENERATION_MANIFEST.json for non-secret source metadata, selected GA4 property ID, report period, and generated source files.\n"
            "- Auction Insights and Wightlink plan inputs remain manual uploads unless explicitly listed otherwise.\n"
        )
        readme_text += source_note
        claude_prompt_text += source_note
        qa_checklist_text += "\n- Check SOURCE_GENERATION_MANIFEST.json when performance/trends came from APIs rather than manual uploads.\n"

    files_manifest = _build_files_manifest(raw_inputs, has_reference_pptx=reference_pptx_bytes is not None, is_monthly=is_monthly)
    if source_generation_payload is not None:
        files_manifest.append(
            {
                "name": "SOURCE_GENERATION_MANIFEST.json",
                "role": "non-secret API source-generation metadata",
                "required": False,
            }
        )
    manifest = {
        "package_type": "claude_handoff",
        "report_family": "Wightlink PPC Monthly" if is_monthly else "Wightlink PPC QBR",
        "client_display_name": "Wightlink",
        "report_mode": report_mode,
        "period_label": period_label,
        "quarter_short": quarter_short,
        "ytd_period_label": ytd_period_label,
        "previous_ytd_period_label": previous_ytd_period_label,
        "generated_at": _format_generated_at(generated_at),
        "reference_deck_url": reference_deck_url,
        "reference_slide_count": REFERENCE_SLIDE_COUNT,
        "target_output_slide_count": streamlit_slide_count if is_monthly else TARGET_OUTPUT_SLIDE_COUNT,
        "streamlit_slide_count": streamlit_slide_count,
        "has_reference_pptx": reference_pptx_bytes is not None,
        "reference_pptx_filename": _reference_pptx_filename(is_monthly) if reference_pptx_bytes is not None else None,
        "has_plan_source": plan_book_csv is not None,
        "plan_source_url": PLAN_SOURCE_URL,
        "headline_kpis": headline_kpis,
        "trend_queries": trend_queries,
        "trend_sources": trend_sources,
        "plan_comparison": {
            "source_table": "PPC Middle Plan Scenario",
            "metrics_available": plan_metrics,
            "missing_months": [],
        },
        "source_sections": [
            {"index": section.index, "title": section.title, "line": section.line, "slide": section.slide}
            for section in source_sections
        ],
        "placeholder_sections": [] if is_monthly else _placeholder_sections(source_sections),
        "excluded_qbr_sections": (
            ["Google Trends", "Auction Insights", "Red Funnel", "Testing", "Other Updates", "Next Steps", "Thank You"]
            if is_monthly
            else []
        ),
        "files": files_manifest,
        "warnings": warnings,
    }
    if source_generation_json is not None:
        manifest["source_generation"] = source_generation_json.get("source_generation", source_generation_json)

    output = BytesIO()
    with ZipFile(output, "w", compression=ZIP_DEFLATED) as archive:
        archive.writestr("report.txt", report_text)
        archive.writestr("original_streamlit_prompt.txt", prompt_text)
        archive.writestr(variables["streamlit_pptx_filename"], generated_pptx_bytes)
        if reference_pptx_bytes is not None:
            archive.writestr(_reference_pptx_filename(is_monthly), reference_pptx_bytes)
        archive.writestr("README_FOR_CLAUDE.txt", readme_text)
        archive.writestr("CLAUDE_PROMPT.txt", claude_prompt_text)
        archive.writestr("SLIDE_MAPPING.csv", slide_mapping_text)
        if not is_monthly:
            archive.writestr("UPDATED_SLIDE_MAPPING_WIGHTLINK_QBR_V2_TEMPLATE.csv", slide_mapping_text)
        archive.writestr("SOURCE_SECTION_INDEX.txt", source_index_text)
        archive.writestr("INPUT_FILES_MANIFEST.txt", input_files_manifest_text)
        archive.writestr("RAW_INPUTS_MANIFEST.txt", input_files_manifest_text)
        archive.writestr("REFERENCE_DECK_OUTLINE.txt", reference_outline_text)
        archive.writestr("QA_CHECKLIST.txt", qa_checklist_text)
        archive.writestr("CHART_QA_ADDENDUM_FOR_CLAUDE.txt", chart_qa_text)
        if source_generation_payload is not None:
            archive.writestr("SOURCE_GENERATION_MANIFEST.json", source_generation_payload)
        if not is_monthly:
            for filename in (
                "GOOGLE_TRENDS_YTD_COMPARISON_RULES.txt",
                "PLAN_COMPARISON_RULES.txt",
                "AUCTION_INSIGHTS_REDFUNNEL_QUARTER_RULES.txt",
                "YTD_MONTHLY_BREAKDOWN_RULES.txt",
                "QA_CHECKLIST_V2_ADDENDUM.txt",
            ):
                archive.writestr(filename, read_asset_text(filename))
        archive.writestr("PACKAGE_MANIFEST.json", json.dumps(manifest, indent=2, ensure_ascii=False))
        for raw_input in raw_inputs:
            archive.writestr(raw_input.archive_name, raw_input.payload)

    return output.getvalue(), manifest


def build_raw_inputs(
    *,
    performance_csv: bytes | Path | str,
    auction_csv: bytes | Path | str | None,
    trend_csv_files: Sequence[bytes | Path | str],
    trend_ytd_current_csv_files: Sequence[bytes | Path | str],
    trend_ytd_previous_csv_files: Sequence[bytes | Path | str],
    red_funnel_auction_csv: bytes | Path | str | None,
    red_funnel_prior_auction_csv: bytes | Path | str | None,
    plan_book_csv: bytes | Path | str | None,
) -> list[WightlinkRawInput]:
    inputs = [
        WightlinkRawInput(
            source_name=_source_name(performance_csv),
            archive_name="source_data/performance.csv",
            role="Primary performance source",
            required=True,
            payload=_read_payload(performance_csv),
        )
    ]

    if auction_csv is not None:
        inputs.append(
            WightlinkRawInput(
                source_name=_source_name(auction_csv),
                archive_name="source_data/auction_insights.csv",
                role="Auction insights source",
                required=True,
                payload=_read_payload(auction_csv),
            )
        )

    if red_funnel_auction_csv is not None:
        inputs.append(
            WightlinkRawInput(
                source_name=_source_name(red_funnel_auction_csv),
                archive_name="source_data/auction_insights_red_funnel_quarter.csv",
                role="Quarter-only Red Funnel Auction Insights source",
                required=True,
                payload=_read_payload(red_funnel_auction_csv),
            )
        )

    if red_funnel_prior_auction_csv is not None:
        inputs.append(
            WightlinkRawInput(
                source_name=_source_name(red_funnel_prior_auction_csv),
                archive_name="source_data/auction_insights_red_funnel_prior_year_quarter.csv",
                role="Same-quarter-prior-year Red Funnel Auction Insights source",
                required=True,
                payload=_read_payload(red_funnel_prior_auction_csv),
            )
        )

    used_trend_names: set[str] = set()
    for trend_file in trend_csv_files:
        payload = _read_payload(trend_file)
        query = extract_trend_query(payload) or _title_from_filename(_source_name(trend_file))
        base_name = f"source_data/google_trends_{_slugify(query)}.csv"
        archive_name = _dedupe_archive_name(base_name, used_trend_names)
        inputs.append(
            WightlinkRawInput(
                source_name=_source_name(trend_file),
                archive_name=archive_name,
                role="Google Trends source",
                required=False,
                payload=payload,
                trend_query=query,
                trend_period="legacy",
            )
        )

    for trend_file in trend_ytd_current_csv_files:
        payload = _read_payload(trend_file)
        query = extract_trend_query(payload) or _title_from_filename(_source_name(trend_file))
        base_name = f"source_data/google_trends_{_slugify(query)}_current_ytd.csv"
        archive_name = _dedupe_archive_name(base_name, used_trend_names)
        inputs.append(
            WightlinkRawInput(
                source_name=_source_name(trend_file),
                archive_name=archive_name,
                role="Current YTD Google Trends source",
                required=False,
                payload=payload,
                trend_query=query,
                trend_period="current_ytd",
            )
        )

    for trend_file in trend_ytd_previous_csv_files:
        payload = _read_payload(trend_file)
        query = extract_trend_query(payload) or _title_from_filename(_source_name(trend_file))
        base_name = f"source_data/google_trends_{_slugify(query)}_previous_ytd.csv"
        archive_name = _dedupe_archive_name(base_name, used_trend_names)
        inputs.append(
            WightlinkRawInput(
                source_name=_source_name(trend_file),
                archive_name=archive_name,
                role="Previous YTD Google Trends source",
                required=False,
                payload=payload,
                trend_query=query,
                trend_period="previous_ytd",
            )
        )

    if plan_book_csv is not None:
        inputs.append(
            WightlinkRawInput(
                source_name=_source_name(plan_book_csv),
                archive_name="source_data/wightlink_plan_2026_27_middle_scenario.csv",
                role="Wightlink plan source",
                required=False,
                payload=_read_payload(plan_book_csv),
            )
        )

    return inputs


def build_monthly_raw_inputs(
    *,
    performance_csv: bytes | Path | str,
    plan_book_csv: bytes | Path | str | None,
) -> list[WightlinkRawInput]:
    inputs = [
        WightlinkRawInput(
            source_name=_source_name(performance_csv),
            archive_name="source_data/performance.csv",
            role="Primary performance source",
            required=True,
            payload=_read_payload(performance_csv),
        )
    ]
    if plan_book_csv is not None:
        suffix = Path(_source_name(plan_book_csv)).suffix or ".csv"
        inputs.append(
            WightlinkRawInput(
                source_name=_source_name(plan_book_csv),
                archive_name=f"source_data/wightlink_plan{suffix}",
                role="Selected-month Wightlink plan source",
                required=False,
                payload=_read_payload(plan_book_csv),
            )
        )
    return inputs


def parse_source_sections(report_text: str) -> list[WightlinkSourceSection]:
    lines = report_text.splitlines()
    sections: list[WightlinkSourceSection] = []
    for index, line in enumerate(lines):
        match = SECTION_RE.match(line.strip())
        if not match:
            continue
        title = match.group("title").strip()
        slide = _extract_slide_number(lines[index + 1 : index + 5])
        sections.append(WightlinkSourceSection(index=len(sections) + 1, title=title, line=index + 1, slide=slide))
    return sections


def extract_period_label(report_text: str) -> str:
    for line in report_text.splitlines():
        if line.startswith("Subtitle:"):
            match = MONTHLY_PERIOD_RE.search(line) or PERIOD_RE.search(line)
            if match:
                return match.group(0)
    match = MONTHLY_PERIOD_RE.search(report_text) or PERIOD_RE.search(report_text)
    return match.group(0) if match else "Unknown period"


def extract_period_short(period_label: str) -> str:
    match = QUARTER_SHORT_RE.search(period_label) or MONTH_SHORT_RE.search(period_label)
    return match.group(0) if match else period_label


def extract_ytd_period_labels(report_text: str, quarter_short: str) -> tuple[str, str]:
    labels: list[str] = []
    for match in re.finditer(r"\bYTD\s+\d{4}\s*\([^)]+\)", report_text):
        label = match.group(0)
        if label not in labels:
            labels.append(label)
    if len(labels) >= 2:
        return labels[0], labels[1]

    month_match = re.search(r"(?P<month>[A-Z][a-z]{2})\s+(?P<year>\d{4})", quarter_short)
    if month_match and not quarter_short.startswith("Q"):
        year = int(month_match.group("year"))
        end_label = month_match.group("month")
        return f"YTD {year} (Jan - {end_label} {year})", f"YTD {year - 1} (Jan - {end_label} {year - 1})"

    quarter_match = re.search(r"Q(?P<quarter>[1-4])\s+(?P<year>\d{4})", quarter_short)
    if not quarter_match:
        return "YTD period unknown", "Previous YTD period unknown"
    quarter = int(quarter_match.group("quarter"))
    year = int(quarter_match.group("year"))
    end_month = quarter * 3
    month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    end_label = month_names[end_month - 1]
    return f"YTD {year} (Jan - {end_label} {year})", f"YTD {year - 1} (Jan - {end_label} {year - 1})"


def extract_headline_kpis(report_text: str) -> dict[str, str]:
    section = _extract_section(report_text, "All Performance Month Summary") or _extract_section(report_text, "All Performance Summary")
    values = {key: "n/a" for key in ("cost", "purchases", "purchase_revenue", "cpa", "roas", "aov")}
    key_map = {
        "Cost": "cost",
        "Purchases": "purchases",
        "Purchase Revenue": "purchase_revenue",
        "CPA": "cpa",
        "ROAS": "roas",
        "AOV": "aov",
    }
    for line in section.splitlines():
        match = KPI_RE.match(line.strip())
        if match:
            values[key_map[match.group("label")]] = match.group("value").strip()
    return values


def extract_plan_metrics(report_text: str) -> list[str]:
    section = _extract_section(report_text, "All Performance Month Summary") or _extract_section(report_text, "All Performance Summary")
    metrics: list[str] = []
    for line in section.splitlines():
        stripped = line.strip()
        match = KPI_RE.match(stripped)
        if match and "Plan:" in stripped:
            metrics.append(match.group("label"))
    return metrics


def build_trend_sources_manifest(raw_inputs: list[WightlinkRawInput]) -> list[dict[str, str]]:
    records: dict[str, dict[str, str]] = {}
    for raw_input in raw_inputs:
        if not raw_input.trend_query:
            continue
        key = _slugify(raw_input.trend_query)
        record = records.setdefault(
            key,
            {
                "standard_name": f"google_trends_{key}",
                "original_filename": "",
                "display_name": raw_input.trend_query,
                "current_ytd_filename": "",
                "previous_ytd_filename": "",
                "normalization_note": "",
            },
        )
        if raw_input.trend_period == "current_ytd":
            record["current_ytd_filename"] = raw_input.archive_name
            record["normalization_note"] = "Google Trends current and previous YTD series were supplied as separate normalized exports."
        elif raw_input.trend_period == "previous_ytd":
            record["previous_ytd_filename"] = raw_input.archive_name
            record["normalization_note"] = "Google Trends current and previous YTD series were supplied as separate normalized exports."
        elif not record["original_filename"]:
            record["original_filename"] = raw_input.archive_name
    return sorted(records.values(), key=lambda record: _trend_sort_key(record["display_name"]))


def _trend_template_variables(trend_sources: list[dict[str, str]]) -> dict[str, str]:
    variables: dict[str, str] = {}
    source_by_term = {_normalise_header(source["display_name"]): source for source in trend_sources}
    for index, term in enumerate(EXPECTED_TREND_TERMS, start=1):
        source = source_by_term.get(_normalise_header(term), {})
        variables[f"trend_{index}_display_name"] = source.get("display_name", term)
        variables[f"trend_{index}_original_filename"] = source.get("original_filename", "not supplied")
        variables[f"trend_{index}_file"] = source.get("original_filename") or source.get("current_ytd_filename") or "not supplied"
        variables[f"trend_{index}_current_ytd_file"] = source.get("current_ytd_filename", "not supplied")
        variables[f"trend_{index}_previous_ytd_file"] = source.get("previous_ytd_filename", "not supplied")
        variables[f"trend_{index}_current_ytd_original_filename"] = _basename_or_not_supplied(source.get("current_ytd_filename"))
        variables[f"trend_{index}_previous_ytd_original_filename"] = _basename_or_not_supplied(source.get("previous_ytd_filename"))
    return variables


def extract_trend_query(payload: bytes) -> str | None:
    text = payload.decode("utf-8-sig", errors="replace")
    for line in text.splitlines():
        if not line.strip():
            continue
        cells = next(csv.reader([line]))
        if len(cells) >= 2 and _normalise_header(cells[0]) in {"time", "date", "week", "month"}:
            query = str(cells[1]).strip()
            return query or None
    return None


def count_pptx_slides(pptx_bytes: bytes) -> int:
    try:
        return len(Presentation(BytesIO(pptx_bytes)).slides)
    except Exception:
        return 0


def find_missing_sections(source_sections: list[WightlinkSourceSection]) -> list[str]:
    titles = {section.title for section in source_sections}
    missing = [section for section in EXPECTED_EXACT_SECTIONS if section not in titles]
    ytd_trend_count = sum(1 for title in titles if title.startswith("Google Trends YTD - "))
    legacy_trend_count = sum(1 for title in titles if title.startswith("Trends - "))
    if ytd_trend_count == 0 and legacy_trend_count < 3:
        missing.append("Three Google Trends YTD detail sections")
    return missing


def find_extra_sections(source_sections: list[WightlinkSourceSection]) -> list[str]:
    extras: list[str] = []
    for section in source_sections:
        if section.title.startswith("Trends - ") or section.title.startswith("Google Trends YTD - "):
            continue
        if section.title not in EXPECTED_EXACT_SECTIONS:
            extras.append(section.title)
    return extras


def find_monthly_missing_sections(source_sections: list[WightlinkSourceSection]) -> list[str]:
    titles = {section.title for section in source_sections}
    required = [
        "Cover",
        "Performance",
        "All Performance Month Summary",
        "All Performance YTD Purchases and Revenue",
        "Brand Month Summary",
        "Brand YTD Purchases and Revenue",
        "Generics Month Summary",
        "Generics YTD Purchases and Revenue",
        "PMax Month Summary",
        "PMax YTD Purchases and Revenue",
    ]
    return [section for section in required if section not in titles]


def find_monthly_extra_sections(source_sections: list[WightlinkSourceSection]) -> list[str]:
    qbr_only_prefixes = ("Google Trends", "Auction Insights")
    qbr_only_titles = {
        "Trends",
        "Agenda",
        "Closing",
        "Thank You",
        "Testing",
        "Other Updates",
        "Next Steps",
        "Other Month Summary",
        "Other YTD Purchases and Revenue",
    }
    extras = []
    for section in source_sections:
        if section.title in qbr_only_titles or section.title.startswith(qbr_only_prefixes):
            extras.append(section.title)
    return extras


def build_source_section_index(
    *,
    report_text: str,
    source_sections: list[WightlinkSourceSection] | None = None,
    period_label: str | None = None,
    missing_sections: list[str] | None = None,
    extra_sections: list[str] | None = None,
    trend_queries: list[str] | None = None,
    ytd_period_label: str | None = None,
    previous_ytd_period_label: str | None = None,
    trend_sources: list[dict[str, str]] | None = None,
    report_mode: str = "quarterly",
) -> str:
    sections = source_sections or parse_source_sections(report_text)
    resolved_period = period_label or extract_period_label(report_text)
    missing = missing_sections if missing_sections is not None else find_missing_sections(sections)
    extra = extra_sections if extra_sections is not None else find_extra_sections(sections)
    queries = trend_queries or []
    sources = trend_sources or []

    output = [
        "SOURCE_SECTION_INDEX.txt",
        "",
        "Client/market: Wightlink",
        f"Period: {resolved_period}",
        f"YTD period: {ytd_period_label or 'not detected'}",
        f"Previous YTD period: {previous_ytd_period_label or 'not detected'}",
        f"Detected source sections: {len(sections)}",
        "",
        "Sections",
    ]
    for section in sections:
        slide = f", slide {section.slide}" if section.slide is not None else ""
        output.append(f"{section.index:02d}. line {section.line}{slide}: [{section.title}]")

    if report_mode == "monthly":
        output.extend(["", "Monthly package note"])
        output.append("- Google Trends and Auction Insights are intentionally excluded from Wightlink monthly packages.")
        output.append("- KPI cards are selected-month values; charts/tables are YTD through the selected month.")
    else:
        output.extend(["", "Google Trends query names from raw inputs"])
        output.extend([f"- {query}" for query in queries] if queries else ["- none supplied"])
        output.extend(["", "Google Trends YTD source pairs"])
        if sources:
            for source in sources:
                output.append(
                    f"- {source.get('display_name')}: current={source.get('current_ytd_filename') or 'not supplied'}; previous={source.get('previous_ytd_filename') or 'not supplied'}"
                )
        else:
            output.append("- none supplied")
        if _has_filename_based_trend_titles(sections):
            output.append("- Note: report trend section titles appear filename-derived; Claude must use the CSV header query names above.")

    output.extend(["", "Missing expected sections"])
    output.extend([f"- {section}" for section in missing] if missing else ["- none"])
    output.extend(["", "Extra sections"])
    output.extend([f"- {section}" for section in extra] if extra else ["- none"])
    output.append("")
    return "\n".join(output)


def build_raw_inputs_manifest(
    *,
    variables: Mapping[str, str],
    raw_inputs: list[WightlinkRawInput],
    missing_required_inputs: list[str],
) -> str:
    if variables.get("report_mode") == "monthly":
        output = [
            "INPUT_FILES_MANIFEST.txt",
            "",
            "Client/market: Wightlink",
            f"Report mode: monthly",
            f"Period: {variables.get('period_label')}",
            "",
            "Source priority",
            "1. report.txt is the source of truth for all values, bullets, tables, MoM, YoY, Plan, and period labels.",
            "2. source_data/performance.csv is the raw performance backup source.",
            "3. source_data/wightlink_plan.* is the raw plan backup source when included.",
            "4. The reference PPTX is visual style only and must not be used for monthly structure.",
            "",
            "Actual raw inputs in this package",
        ]
        for raw_input in raw_inputs:
            output.append(f"- {raw_input.archive_name}: {raw_input.role}; source file {raw_input.source_name}")
        output.extend(["", "Missing raw inputs"])
        output.extend([f"- {name}" for name in missing_required_inputs] if missing_required_inputs else ["- none"])
        output.append("")
        return "\n".join(output)

    output = [render_asset_template("INPUT_FILES_MANIFEST_TEMPLATE.txt", variables).rstrip(), "", "Actual raw inputs in this package"]
    for raw_input in raw_inputs:
        line = f"- {raw_input.archive_name}: {raw_input.role}; source file {raw_input.source_name}"
        if raw_input.trend_query:
            line += f"; query header {raw_input.trend_query}"
        output.append(line)
    output.extend(["", "Missing raw inputs"])
    output.extend([f"- {name}" for name in missing_required_inputs] if missing_required_inputs else ["- none"])
    output.append("")
    return "\n".join(output)


def build_monthly_readme_text(variables: Mapping[str, str]) -> str:
    return f"""
Wightlink {variables['quarter_short']} Monthly PPC Report Deck Handoff

Goal
Transform the Streamlit output in this upload pack into a polished {variables['target_output_slide_count']}-slide monthly Wightlink PPC report deck. The monthly structure is defined only by SLIDE_MAPPING.csv and report.txt.

Files in this pack
- report.txt: Source of truth for all {variables['period_label']} values, MoM/YoY/Plan card comparisons, YTD chart tables, bullets, and period labels.
- {variables['streamlit_pptx_filename']}: Streamlit-generated PowerPoint. Use it to understand structure and chart coverage.
- {variables['reference_pptx_filename']}: QBR deck export for visual style only, if included. It is not the monthly slide map. Do not copy its slide order or QBR-only sections.
- original_streamlit_prompt.txt: Monthly Streamlit prompt for background only.
- SLIDE_MAPPING.csv: Monthly source-section execution map.
- SOURCE_SECTION_INDEX.txt: Index of report.txt sections and line numbers.
- INPUT_FILES_MANIFEST.txt and RAW_INPUTS_MANIFEST.txt: Raw input roles.
- QA_CHECKLIST.txt: Required checks before returning the completed deck.
- CHART_QA_ADDENDUM_FOR_CLAUDE.txt: Required chart rendering and visual QA rules.
- REFERENCE_DECK_OUTLINE.txt: Visual-reference note for the Wightlink QBR deck.
- PACKAGE_MANIFEST.json: Package metadata.
- source_data/: Raw uploaded inputs. Monthly packages should contain the performance CSV and the plan workbook/CSV if uploaded.

Monthly rules
- Target deck is {variables['target_output_slide_count']} slides.
- Use SLIDE_MAPPING.csv as the only target slide structure.
- KPI cards show the selected month only.
- KPI cards must include MoM, YoY, and Plan where report.txt provides them.
- Spend/Cost MoM and YoY comparison text is always neutral grey, never red or green.
- Plan comparisons are selected-month actuals versus selected-month plan.
- Tables and charts show YTD through the selected month.
- YTD tables belong on the preceding Month Summary slides.
- Each YTD Purchases and Revenue slide must contain two YoY charts: Purchases current YTD vs prior-year YTD, and Revenue current YTD vs prior-year YTD.
- Do not include the old combined Purchases-and-Revenue chart.
- Do not include monthly campaign Other slides; monthly Wightlink includes All Performance, Brand, Generics, PMax, Ferry, and Routes where available.
- Do not add Google Trends, Auction Insights, Red Funnel, Testing, Other Updates, Next Steps, Thank You, or other QBR-only slides unless those sections are present in SLIDE_MAPPING.csv.
- Use report.txt values exactly. Do not recalculate or round differently unless resolving a source inconsistency.
- Remove all old Q1/Q2/quarterly wording from any reused visual reference objects.

Recommended workflow
1. Read report.txt and use SLIDE_MAPPING.csv as the execution map.
2. Use the Streamlit PPTX to understand monthly/YTD source coverage.
3. Use {variables['reference_pptx_filename']} only for visual style. Do not use it for structure.
4. Generate chart PNGs from report.txt data where needed.
5. Apply CHART_QA_ADDENDUM_FOR_CLAUDE.txt to every chart slide and inspect full-slide screenshots or thumbnails.
6. Run QA_CHECKLIST.txt before returning the finished deck.
""".strip()


def build_monthly_claude_prompt_text(variables: Mapping[str, str]) -> str:
    return f"""
You are preparing a Wightlink {variables['quarter_short']} monthly PPC report deck.

I am uploading:
- report.txt
- {variables['streamlit_pptx_filename']}
- {variables['reference_pptx_filename']}, if included
- original_streamlit_prompt.txt
- SLIDE_MAPPING.csv
- SOURCE_SECTION_INDEX.txt
- INPUT_FILES_MANIFEST.txt
- QA_CHECKLIST.txt
- CHART_QA_ADDENDUM_FOR_CLAUDE.txt
- PACKAGE_MANIFEST.json
- source_data/ raw performance and plan inputs

Task
Create a {variables['target_output_slide_count']}-slide monthly report deck populated with {variables['period_label']} Wightlink data from report.txt.

Source priority
1. report.txt is the source of truth for all numbers, tables, bullets, dates, MoM, YoY, Plan, and YTD values.
2. SLIDE_MAPPING.csv defines the only slide structure to build.
3. {variables['streamlit_pptx_filename']} shows the intermediate Streamlit structure and charts.
4. {variables['reference_pptx_filename']}, if included, defines visual style only. Do not use its slide order or QBR-only sections.
5. source_data/ files are audit/backup sources.
6. original_streamlit_prompt.txt is background only.

Hard requirements
- Build exactly {variables['target_output_slide_count']} monthly slides unless SLIDE_MAPPING.csv says otherwise.
- Do not copy the quarterly/QBR deck structure.
- Do not add Google Trends, Auction Insights, Red Funnel, Testing, Other Updates, Next Steps, Thank You, or other QBR-only slides unless present in SLIDE_MAPPING.csv.
- KPI cards show the selected month only.
- KPI cards must show MoM, YoY, and Plan where report.txt provides them.
- Spend/Cost MoM and YoY comparison text must be neutral grey, never red or green.
- Plan comparisons are selected-month actuals versus selected-month plan.
- Tables and charts show YTD through the selected month.
- Put YTD tables on the preceding Month Summary slides.
- Build each YTD Purchases and Revenue slide with two YoY charts: Purchases current YTD vs prior-year YTD, and Revenue current YTD vs prior-year YTD.
- Do not include the old combined Purchases-and-Revenue chart.
- Include every Wightlink scope present in report.txt: All Performance, Brand, Generics, PMax, Ferry, and Routes where available.
- Do not add monthly campaign Other slides.
- Preserve the Wightlink/Summon visual system: dark title/section slides, Wightlink accent colors, KPI cards, tables, footers, typography, chart placement, and spacing.
- Remove all old Q1/Q2/quarterly wording from reused style objects.
- Apply CHART_QA_ADDENDUM_FOR_CLAUDE.txt to every chart slide. Charts must not be squashed, clipped, or overlapping.
- Run QA_CHECKLIST.txt before returning.

Output
Return the finished Google Slides deck URL if you edited Google Slides directly. If you used the PPTX bridge, return the finished PPTX only and note that it should be imported into Google Slides.
""".strip()


def build_monthly_slide_mapping_text(source_sections: list[WightlinkSourceSection]) -> str:
    output = StringIO()
    fieldnames = ["target_slide", "reference_title", "source_section_or_action", "required_action"]
    writer = csv.DictWriter(output, fieldnames=fieldnames)
    writer.writeheader()
    for section in source_sections:
        writer.writerow(
            {
                "target_slide": section.slide or section.index,
                "reference_title": section.title,
                "source_section_or_action": section.title,
                "required_action": _monthly_slide_action(section.title),
            }
        )
    return output.getvalue()


def build_monthly_qa_checklist_text(variables: Mapping[str, str]) -> str:
    return f"""
Required QA before delivery

Deck identity
- Final deck is a {variables['target_output_slide_count']}-slide monthly Wightlink PPC report for {variables['period_label']}.
- Final deck follows SLIDE_MAPPING.csv, not the quarterly/QBR reference deck.
- Client/market naming is Wightlink.

Stale content
- No Google Trends slides remain unless present in SLIDE_MAPPING.csv.
- No Auction Insights slides remain unless present in SLIDE_MAPPING.csv.
- No Red Funnel slides remain unless present in SLIDE_MAPPING.csv.
- No Testing, Other Updates, Next Steps, or Thank You slides remain unless present in SLIDE_MAPPING.csv.
- No Q1/Q2/quarter/quarterly date labels remain.
- No "Wightlink QBR" title remains unless it is explicitly used as a visual-reference note.
- No old QBR slide numbers are used for chart QA or section mapping.

Data integrity
- KPI values match report.txt selected-month KPI card rows exactly.
- MoM, YoY, and Plan values match report.txt exactly.
- KPI cards use selected-month values only.
- Plan comparisons are selected-month actuals versus selected-month plan.
- Tables and charts use YTD data through the selected month.
- Spend/Cost MoM and YoY comparison text is neutral grey, never red or green.
- YTD tables are on the preceding Month Summary slides.
- YTD Purchases and Revenue slides contain two YoY charts: Purchases and Revenue.
- No old combined Purchases-and-Revenue chart remains.
- No monthly campaign Other slides remain.
- Every report.txt source section is represented in the deck.

Visual fidelity
- Typography, colors, title bands, section dividers, KPI card styling, tables, footers, and chart treatments match the Wightlink/Summon reference style.
- Text does not overflow or wrap awkwardly in title/footer bands.
- Tables fit within slide bounds and remain readable.
- Every chart slide has been checked against CHART_QA_ADDENDUM_FOR_CLAUDE.txt.
- No chart labels, value annotations, legends, or axis labels are clipped or overlapping.
- Chart images are placed with fit/contain behavior and are not stretched non-proportionally.
- No objects overlap unintentionally.

Final response
- If the output is native Google Slides, return the new Google Slides URL only.
- If the output is PPTX, return the finished PPTX and state it is ready to import into Google Slides.
""".strip()


def build_monthly_chart_qa_addendum_text(source_sections: list[WightlinkSourceSection], variables: Mapping[str, str]) -> str:
    chart_sections = [section for section in source_sections if "YTD Purchases and Revenue" in section.title]
    summary_sections = [section for section in source_sections if "Month Summary" in section.title]

    def _format_sections(sections: list[WightlinkSourceSection]) -> str:
        if not sections:
            return "- none"
        return "\n".join(f"- Slide {section.slide or section.index}: {section.title}" for section in sections)

    return f"""
Monthly Chart QA Addendum for Claude

Problem to fix
The monthly deck content can be correct while chart images are too tightly rendered or incorrectly scaled inside the slide slots. This causes labels, value annotations, legends, and axis text to look squashed, clipped, or overlapping.

Target structure
- This is a {variables['target_output_slide_count']}-slide monthly deck for {variables['period_label']}.
- Use SLIDE_MAPPING.csv for slide numbers.
- Do not use old quarterly/QBR slide numbers.
- {variables['reference_pptx_filename']} is a visual style reference only, not a structural map.

YTD purchases and revenue YoY chart slides to inspect
{_format_sections(chart_sections)}

Summary slides with KPI cards to inspect
{_format_sections(summary_sections)}

Chart rendering rules
- Preserve the intended slide slot size and position from the monthly output or approved visual style.
- Do not stretch chart images non-proportionally.
- Use fit/contain placement, not crop/fill placement.
- Export chart images at high resolution, ideally 2x or 3x the displayed slide size.
- Leave internal padding inside every exported chart image so labels cannot be cut off after insertion.
- Add at least 10-15% internal padding around chart labels.
- Avoid chart titles inside the image when the slide already has a title.
- Keep legends readable but compact; legends should not push charts smaller than the intended chart area.
- For bar charts, make sure value labels are not clipped at the top or right edge. Increase the axis maximum or add margin.
- Keep the Wightlink/Summon monthly color system and chart frame treatment.
- Each YTD Purchases and Revenue slide must contain two readable charts: Purchases YoY and Revenue YoY.

Required visual QA
For every chart slide, create or inspect a full-slide screenshot/thumbnail after the chart replacement. The pass is not complete until each chart passes these checks:
- No number is clipped.
- No label overlaps another label, legend, bar, or axis.
- No chart appears squeezed vertically or horizontally.
- The chart is centered in its intended slot.
- The chart has professional whitespace.
- The chart does not collide with bullets, table text, headers, footers, or logos.
- The slide still reads as a monthly Wightlink report slide.

Final response
List the slides whose charts were changed and confirm that screenshots/thumbnails were checked after the fix.
""".strip()


def build_monthly_reference_outline_text(variables: Mapping[str, str]) -> str:
    return f"""
REFERENCE_DECK_OUTLINE.txt

Monthly package note
- {variables['reference_pptx_filename']} is included only as a Wightlink/Summon visual style reference.
- Do not copy the QBR slide order.
- Do not add QBR-only Google Trends, Auction Insights, Red Funnel, Testing, Other Updates, Next Steps, or Thank You slides.
- Use SLIDE_MAPPING.csv as the monthly slide outline.
- Use report.txt as the source of truth for {variables['period_label']} content.
""".strip()


def _monthly_slide_action(title: str) -> str:
    if title == "Cover":
        return "Update title, client, and period. Make clear this is a monthly Wightlink PPC report."
    if title == "Performance":
        return "Keep as a divider or section header. Use monthly period labeling."
    if "Month Summary" in title:
        return "Update selected-month KPI cards from report.txt. Include MoM, YoY, and Plan lines where provided, keep Spend/Cost MoM and YoY text neutral grey, and place the YTD table on this slide. Preserve bullets exactly."
    if "YTD Purchases and Revenue" in title:
        return "Create or replace two YoY charts from report.txt: Purchases current YTD vs prior-year YTD and Revenue current YTD vs prior-year YTD. Do not use the old combined purchases/revenue chart. Apply CHART_QA_ADDENDUM_FOR_CLAUDE.txt."
    if title == "Ferry & Routes":
        return "Keep as a divider if Ferry/Routes sections exist in report.txt."
    return "Represent this report.txt section exactly; do not add QBR-only content."


def render_asset_template(filename: str, variables: Mapping[str, str]) -> str:
    text = read_asset_text(filename)
    for key, value in variables.items():
        text = text.replace(f"{{{{{key}}}}}", str(value))
    return text


def read_asset_text(filename: str) -> str:
    return (ASSETS_DIR / filename).read_text(encoding="utf-8")


def _extract_section(report_text: str, title: str) -> str:
    lines = report_text.splitlines()
    start_index: int | None = None
    for index, line in enumerate(lines):
        if line.strip() == f"[{title}]":
            start_index = index
            break
    if start_index is None:
        return ""
    end_index = len(lines)
    for index in range(start_index + 1, len(lines)):
        if SECTION_RE.match(lines[index].strip()):
            end_index = index
            break
    return "\n".join(lines[start_index:end_index])


def _extract_slide_number(lines: Sequence[str]) -> int | None:
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("Slide:"):
            value = stripped.split(":", 1)[1].strip()
            try:
                return int(value)
            except ValueError:
                return None
    return None


def _max_report_slide_number(source_sections: list[WightlinkSourceSection]) -> int:
    slides = [section.slide for section in source_sections if section.slide is not None]
    return max(slides) if slides else 0


def _has_filename_based_trend_titles(source_sections: list[WightlinkSourceSection]) -> bool:
    for section in source_sections:
        if not section.title.startswith("Trends - "):
            continue
        trend_title = section.title.removeprefix("Trends - ").strip()
        if re.search(r"\b\d{8}\b", trend_title) or trend_title.lower().endswith(".csv"):
            return True
    return False


def _build_v2_source_warnings(
    trend_sources: list[dict[str, str]],
    raw_inputs: list[WightlinkRawInput],
    plan_book_csv: bytes | Path | str | None,
) -> list[str]:
    warnings: list[str] = []
    source_by_term = {_normalise_header(source["display_name"]): source for source in trend_sources}
    for term in EXPECTED_TREND_TERMS:
        source = source_by_term.get(_normalise_header(term))
        if not source or not source.get("current_ytd_filename") or not source.get("previous_ytd_filename"):
            warnings.append(f"Missing current/previous YTD Google Trends pair for {term}.")
    if not any(raw_input.archive_name == "source_data/auction_insights_red_funnel_quarter.csv" for raw_input in raw_inputs):
        warnings.append("Missing quarter-only Red Funnel Auction Insights source.")
    if not any(raw_input.archive_name == "source_data/auction_insights_red_funnel_prior_year_quarter.csv" for raw_input in raw_inputs):
        warnings.append("Missing same-quarter-prior-year Red Funnel Auction Insights source; Red Funnel YoY change column may be unavailable.")
    if plan_book_csv is None:
        warnings.append("Missing Wightlink plan source; KPI Plan lines may be unavailable.")
    return warnings


def _placeholder_sections(source_sections: list[WightlinkSourceSection]) -> list[str]:
    titles = {section.title for section in source_sections}
    placeholders = []
    for title in ("Testing", "Homepage A/B Test Results", "Microsoft Ads - Performance Max Test", "Next Steps", "Other Updates", "Ad Monitor Now Live"):
        if title not in titles:
            placeholders.append(title)
    if not any(title.startswith("Top Performing Campaigns") for title in titles):
        placeholders.append("Top Performing Campaigns excl. Brand")
    return placeholders


def _reference_pptx_filename(is_monthly: bool) -> str:
    return MONTHLY_REFERENCE_PPTX_FILENAME if is_monthly else DEFAULT_REFERENCE_PPTX_PATH.name


def _missing_required_raw_inputs(raw_inputs: list[WightlinkRawInput], *, is_monthly: bool = False) -> list[str]:
    names = {raw_input.archive_name for raw_input in raw_inputs}
    required = ["source_data/performance.csv"]
    if not is_monthly:
        required.append("source_data/auction_insights.csv")
    return [name for name in required if name not in names]


def _build_files_manifest(raw_inputs: list[WightlinkRawInput], *, has_reference_pptx: bool, is_monthly: bool = False) -> list[dict]:
    if is_monthly:
        files = [
            {"name": "report.txt", "role": "source report text", "required": True},
            {"name": "original_streamlit_prompt.txt", "role": "original Streamlit prompt for background only", "required": True},
            {"name": "wightlink_monthly_streamlit_output.pptx", "role": "Streamlit-generated monthly PPTX intermediate source", "required": True},
            {
                "name": MONTHLY_REFERENCE_PPTX_FILENAME,
                "role": "PPTX export of the Wightlink QBR visual reference deck; style only, not monthly structure",
                "required": False,
            },
            {"name": "README_FOR_CLAUDE.txt", "role": "operator README", "required": True},
            {"name": "CLAUDE_PROMPT.txt", "role": "Claude execution prompt", "required": True},
            {"name": "SLIDE_MAPPING.csv", "role": "monthly report slide mapping", "required": True},
            {"name": "SOURCE_SECTION_INDEX.txt", "role": "report section line index", "required": True},
            {"name": "INPUT_FILES_MANIFEST.txt", "role": "raw input descriptions", "required": True},
            {"name": "RAW_INPUTS_MANIFEST.txt", "role": "raw input descriptions alias", "required": True},
            {"name": "REFERENCE_DECK_OUTLINE.txt", "role": "visual-reference-only note", "required": True},
            {"name": "QA_CHECKLIST.txt", "role": "required final deck QA checklist", "required": True},
            {"name": "CHART_QA_ADDENDUM_FOR_CLAUDE.txt", "role": "required chart rendering and screenshot QA rules", "required": True},
            {"name": "PACKAGE_MANIFEST.json", "role": "machine-readable package metadata", "required": True},
        ]
        if not has_reference_pptx:
            files = [file for file in files if file["name"] != MONTHLY_REFERENCE_PPTX_FILENAME]
        for raw_input in raw_inputs:
            files.append({"name": raw_input.archive_name, "role": raw_input.role, "required": raw_input.required})
        return files

    files = [
        {"name": "report.txt", "role": "source report text", "required": True},
        {"name": "original_streamlit_prompt.txt", "role": "original Streamlit prompt for background only", "required": True},
        {"name": "wightlink_streamlit_output.pptx", "role": "Streamlit-generated PPTX intermediate source", "required": True},
        {
            "name": "reference_deck_exported_from_google_slides.pptx",
            "role": "PPTX export of the Wightlink Google Slides visual reference deck",
            "required": False,
        },
        {"name": "README_FOR_CLAUDE.txt", "role": "operator README", "required": True},
        {"name": "CLAUDE_PROMPT.txt", "role": "Claude execution prompt", "required": True},
        {"name": "SLIDE_MAPPING.csv", "role": "22-slide lean V2 reference deck mapping", "required": True},
        {"name": "UPDATED_SLIDE_MAPPING_WIGHTLINK_QBR_V2_TEMPLATE.csv", "role": "22-slide lean V2 reference deck mapping alias", "required": True},
        {"name": "SOURCE_SECTION_INDEX.txt", "role": "report section line index", "required": True},
        {"name": "INPUT_FILES_MANIFEST.txt", "role": "raw input descriptions and query names", "required": True},
        {"name": "RAW_INPUTS_MANIFEST.txt", "role": "raw input descriptions and query names alias", "required": True},
        {"name": "REFERENCE_DECK_OUTLINE.txt", "role": "Wightlink reference deck outline", "required": True},
        {"name": "QA_CHECKLIST.txt", "role": "required final deck QA checklist", "required": True},
        {"name": "CHART_QA_ADDENDUM_FOR_CLAUDE.txt", "role": "required chart rendering and screenshot QA rules", "required": True},
        {"name": "GOOGLE_TRENDS_YTD_COMPARISON_RULES.txt", "role": "V2 Google Trends comparison rules", "required": True},
        {"name": "PLAN_COMPARISON_RULES.txt", "role": "V2 plan comparison rules", "required": True},
        {"name": "AUCTION_INSIGHTS_REDFUNNEL_QUARTER_RULES.txt", "role": "V2 Red Funnel quarter auction rules", "required": True},
        {"name": "YTD_MONTHLY_BREAKDOWN_RULES.txt", "role": "V2 YTD monthly breakdown rules", "required": True},
        {"name": "QA_CHECKLIST_V2_ADDENDUM.txt", "role": "V2 QA addendum", "required": True},
        {"name": "PACKAGE_MANIFEST.json", "role": "machine-readable package metadata", "required": True},
    ]
    if not has_reference_pptx:
        files = [file for file in files if file["name"] != "reference_deck_exported_from_google_slides.pptx"]
    for raw_input in raw_inputs:
        files.append({"name": raw_input.archive_name, "role": raw_input.role, "required": raw_input.required})
    return files


def _read_payload(payload: bytes | Path | str) -> bytes:
    if isinstance(payload, bytes):
        return payload
    return Path(payload).read_bytes()


def _read_optional_payload(payload: bytes | Path | str | None) -> bytes | None:
    if payload is None:
        return None
    if isinstance(payload, bytes):
        return payload
    path = Path(payload)
    if not path.exists():
        return None
    return path.read_bytes()


def _read_optional_json_manifest(payload: Mapping[str, Any] | bytes | Path | str | None) -> tuple[bytes | None, dict[str, Any] | None]:
    if payload is None:
        return None, None
    if isinstance(payload, Mapping):
        json_payload = dict(payload)
        return json.dumps(json_payload, indent=2, ensure_ascii=False).encode("utf-8"), json_payload
    raw_payload = _read_optional_payload(payload)
    if raw_payload is None:
        return None, None
    try:
        parsed = json.loads(raw_payload.decode("utf-8"))
    except Exception:
        parsed = None
    return raw_payload, parsed if isinstance(parsed, dict) else None


def _source_name(payload: bytes | Path | str) -> str:
    if isinstance(payload, bytes):
        return "uploaded_bytes.csv"
    return Path(payload).name


def _dedupe_archive_name(base_name: str, used_names: set[str]) -> str:
    if base_name not in used_names:
        used_names.add(base_name)
        return base_name
    stem = base_name.removesuffix(".csv")
    counter = 2
    while f"{stem}_{counter}.csv" in used_names:
        counter += 1
    archive_name = f"{stem}_{counter}.csv"
    used_names.add(archive_name)
    return archive_name


def _format_generated_at(value: datetime | None) -> str:
    generated_at = value or datetime.now(timezone.utc)
    if generated_at.tzinfo is None:
        generated_at = generated_at.replace(tzinfo=timezone.utc)
    return generated_at.astimezone(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def _normalise_header(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", value.strip().lower())


def _slugify(value: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "_", value.strip().lower()).strip("_")
    return slug or "trend"


def _title_from_filename(filename: str) -> str:
    return Path(filename).stem.replace("_", " ").replace("-", " ").title()


def _trend_sort_key(value: str) -> tuple[int, str]:
    normalized = _normalise_header(value)
    for index, term in enumerate(EXPECTED_TREND_TERMS):
        if normalized == _normalise_header(term):
            return index, normalized
    return len(EXPECTED_TREND_TERMS), normalized


def _basename_or_not_supplied(value: str | None) -> str:
    if not value:
        return "not supplied"
    return Path(value).name
