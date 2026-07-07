from __future__ import annotations

import os
import math
from pathlib import Path
from typing import Any

os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib")
import matplotlib

matplotlib.use("Agg", force=True)
import matplotlib.pyplot as plt
import pandas as pd
from matplotlib.ticker import FuncFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class WightlinkPptxBuilder:
    def __init__(self, output_path: str | Path, charts_dir: str | Path) -> None:
        self.output_path = Path(output_path)
        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        self.charts_dir = Path(charts_dir)
        self.charts_dir.mkdir(parents=True, exist_ok=True)
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.bg = RGBColor(247, 249, 252)
        self.text_primary = RGBColor(20, 42, 77)
        self.text_secondary = RGBColor(90, 90, 90)
        self.text_body = RGBColor(45, 45, 45)
        self.accent = RGBColor(12, 84, 96)
        self.light_fill = RGBColor(232, 239, 247)

    def build(self, slides: list[dict[str, Any]]) -> Path:
        for slide in slides:
            self._render_slide(slide)
        self.prs.save(str(self.output_path))
        return self.output_path

    def build_trend_chart(self, trend_section: dict[str, Any], filename: str) -> Path:
        output = self.charts_dir / filename
        labels = trend_section.get("labels", [])
        series = trend_section.get("series", [])
        if not labels or not series:
            return self._plot_empty(output, "No trend data")
        x_values = list(range(len(labels)))
        tick_step = max(1, math.ceil(len(labels) / 8))
        tick_positions = x_values[::tick_step]
        if tick_positions[-1] != x_values[-1]:
            tick_positions.append(x_values[-1])

        fig, ax = plt.subplots(figsize=(7.2, 3.8))
        palette = ["#D63C31", "#A3A3A3"] if trend_section.get("chart_style") == "ytd_comparison" else ["#0C5460", "#94A3B8", "#1D4ED8", "#14B8A6"]
        for index, item in enumerate(series):
            data = item.get("data", [])
            marker = "o" if len(data) <= 20 else None
            ax.plot(
                x_values,
                data,
                linewidth=2.3,
                marker=marker,
                markersize=4.5,
                label=item.get("name", f"Series {index+1}"),
                color=item.get("color") or palette[index % len(palette)],
            )
        ax.set_title(trend_section.get("title", "Google Trends"), fontsize=15)
        ax.set_xticks(tick_positions, [labels[index] for index in tick_positions])
        ax.tick_params(axis="x", rotation=35, labelsize=8)
        ax.tick_params(axis="y", labelsize=9)
        ax.grid(axis="y", alpha=0.2)
        ax.margins(x=0.02)
        if len(series) > 1:
            ax.legend(fontsize=9, loc="upper left")
        plt.tight_layout()
        fig.savefig(output, dpi=180)
        plt.close(fig)
        return output

    def build_performance_chart(self, scope: dict[str, Any], filename: str, left_metric: str, right_metric: str | None, title: str) -> Path:
        output = self.charts_dir / filename
        monthly = scope.get("monthly", [])
        if not monthly:
            return self._plot_empty(output, "No performance data")
        df = pd.DataFrame(monthly)
        labels = df["month_label"].tolist()
        fig, ax1 = plt.subplots(figsize=(7.8, 4.0))
        ax1.plot(labels, df[left_metric], color="#0C5460", linewidth=2.3, marker="o", label=left_metric.upper())
        ax1.set_title(title, fontsize=14)
        ax1.tick_params(axis="both", labelsize=9)
        ax1.grid(axis="y", alpha=0.2)
        lines, labels_accum = ax1.get_legend_handles_labels()
        if right_metric and right_metric in df.columns and df[right_metric].notna().any():
            ax2 = ax1.twinx()
            ax2.plot(labels, df[right_metric], color="#1D4ED8", linewidth=2.0, marker="o", label=right_metric.upper())
            ax2.tick_params(axis="y", labelsize=9)
            right_lines, right_labels = ax2.get_legend_handles_labels()
            lines += right_lines
            labels_accum += right_labels
        ax1.legend(lines, labels_accum, fontsize=9, loc="upper left")
        plt.tight_layout()
        fig.savefig(output, dpi=180)
        plt.close(fig)
        return output

    def build_yoy_performance_chart(
        self,
        current_scope: dict[str, Any],
        prior_scope: dict[str, Any] | None,
        filename: str,
        left_metric: str,
        right_metric: str | None,
        title: str,
        current_label: str,
        prior_label: str,
    ) -> Path:
        output = self.charts_dir / filename
        current_monthly = current_scope.get("monthly", [])
        if not current_monthly:
            return self._plot_empty(output, "No performance data")

        current_df = pd.DataFrame(current_monthly)
        labels = current_df["month_label"].tolist()
        fig, ax1 = plt.subplots(figsize=(7.8, 4.0))

        ax1.plot(labels, current_df[left_metric], color="#0C5460", linewidth=2.3, marker="o", label=f"{current_label} {left_metric.upper()}")

        if prior_scope and prior_scope.get("monthly"):
            prior_df = pd.DataFrame(prior_scope["monthly"])
            prior_df = prior_df.set_index("month_label").reindex(labels).reset_index()
            ax1.plot(
                labels,
                prior_df[left_metric],
                color="#94A3B8",
                linewidth=2.1,
                linestyle="--",
                marker="o",
                label=f"{prior_label} {left_metric.upper()}",
            )

        ax1.set_title(title, fontsize=14)
        ax1.tick_params(axis="both", labelsize=9)
        ax1.grid(axis="y", alpha=0.2)
        lines, labels_accum = ax1.get_legend_handles_labels()

        if right_metric and right_metric in current_df.columns and current_df[right_metric].notna().any():
            ax2 = ax1.twinx()
            ax2.plot(
                labels,
                current_df[right_metric],
                color="#1D4ED8",
                linewidth=2.0,
                marker="o",
                label=f"{current_label} {right_metric.upper()}",
            )
            if prior_scope and prior_scope.get("monthly"):
                prior_df = pd.DataFrame(prior_scope["monthly"])
                prior_df = prior_df.set_index("month_label").reindex(labels).reset_index()
                if right_metric in prior_df.columns and prior_df[right_metric].notna().any():
                    ax2.plot(
                        labels,
                        prior_df[right_metric],
                        color="#60A5FA",
                        linewidth=1.9,
                        linestyle="--",
                        marker="o",
                        label=f"{prior_label} {right_metric.upper()}",
                    )
            ax2.tick_params(axis="y", labelsize=9)
            right_lines, right_labels = ax2.get_legend_handles_labels()
            lines += right_lines
            labels_accum += right_labels

        ax1.legend(lines, labels_accum, fontsize=8, loc="upper left")
        plt.tight_layout()
        fig.savefig(output, dpi=180)
        plt.close(fig)
        return output

    def build_yoy_bar_chart(
        self,
        monthly_rows: list[dict[str, Any]],
        filename: str,
        current_metric: str,
        prior_metric: str,
        title: str,
        current_label: str,
        prior_label: str,
    ) -> Path:
        output = self.charts_dir / filename
        if not monthly_rows:
            return self._plot_empty(output, "No YoY data")

        df = pd.DataFrame(monthly_rows)
        if df.empty or "month_label" not in df.columns:
            return self._plot_empty(output, "No YoY data")

        labels = df["month_label"].tolist()
        x_positions = range(len(labels))
        width = 0.34
        current_values = df[current_metric].fillna(0) if current_metric in df.columns else [0] * len(labels)
        prior_values = df[prior_metric].fillna(0) if prior_metric in df.columns else [0] * len(labels)

        fig, ax = plt.subplots(figsize=(8.4, 4.4))
        current_bars = ax.bar(
            [position - width / 2 for position in x_positions],
            current_values,
            width=width,
            color="#C7372F",
            label=current_label,
        )
        prior_bars = ax.bar(
            [position + width / 2 for position in x_positions],
            prior_values,
            width=width,
            color="#B8B5B1",
            label=prior_label,
        )
        ax.set_title(title.upper(), fontsize=12, color="#888888", fontweight="bold")
        ax.set_xticks(list(x_positions), labels)
        ax.yaxis.set_major_formatter(FuncFormatter(lambda value, _: f"{value:,.0f}"))
        ax.tick_params(axis="both", labelsize=9, colors="#888888")
        ax.grid(axis="y", alpha=0.22)
        ax.spines[["top", "right", "left"]].set_visible(False)
        ax.legend(fontsize=9, loc="upper center", bbox_to_anchor=(0.5, -0.08), ncol=2, frameon=False)
        for bars in (current_bars, prior_bars):
            ax.bar_label(bars, labels=[f"{bar.get_height():,.0f}" if bar.get_height() else "" for bar in bars], padding=3, fontsize=8)
        plt.tight_layout()
        fig.savefig(output, dpi=180)
        plt.close(fig)
        return output

    def build_monthly_purchases_revenue_chart(self, scope: dict[str, Any], filename: str, title: str) -> Path:
        output = self.charts_dir / filename
        monthly = scope.get("monthly", [])
        if not monthly:
            return self._plot_empty(output, "No monthly data")

        df = pd.DataFrame(monthly)
        if df.empty or "month_label" not in df.columns:
            return self._plot_empty(output, "No monthly data")

        labels = df["month_label"].tolist()
        x_positions = range(len(labels))
        width = 0.34
        purchases = df["purchases"].fillna(0) if "purchases" in df.columns else [0] * len(labels)
        revenue = df["purchase_revenue"].fillna(0) if "purchase_revenue" in df.columns else [0] * len(labels)

        fig, ax1 = plt.subplots(figsize=(8.4, 4.4))
        ax2 = ax1.twinx()
        purchase_bars = ax1.bar(
            [position - width / 2 for position in x_positions],
            purchases,
            width=width,
            color="#C7372F",
            label="Purchases",
        )
        revenue_bars = ax2.bar(
            [position + width / 2 for position in x_positions],
            revenue,
            width=width,
            color="#B8B5B1",
            label="Revenue",
        )
        ax1.set_title(title.upper(), fontsize=12, color="#888888", fontweight="bold")
        ax1.set_xticks(list(x_positions), labels)
        ax1.yaxis.set_major_formatter(FuncFormatter(lambda value, _: f"{value:,.0f}"))
        ax2.yaxis.set_major_formatter(FuncFormatter(_compact_currency_tick))
        ax1.tick_params(axis="both", labelsize=9, colors="#888888")
        ax2.tick_params(axis="y", labelsize=9, colors="#888888")
        ax1.grid(axis="y", alpha=0.22)
        ax1.spines[["top", "right", "left"]].set_visible(False)
        ax2.spines[["top", "right", "left"]].set_visible(False)
        ax1.legend([purchase_bars, revenue_bars], ["Purchases", "Revenue"], fontsize=9, loc="upper center", bbox_to_anchor=(0.5, -0.08), ncol=2, frameon=False)
        ax1.bar_label(purchase_bars, labels=[f"{bar.get_height():,.0f}" if bar.get_height() else "" for bar in purchase_bars], padding=3, fontsize=8)
        ax2.bar_label(revenue_bars, labels=[_compact_currency_label(bar.get_height()) if bar.get_height() else "" for bar in revenue_bars], padding=3, fontsize=8)
        plt.tight_layout()
        fig.savefig(output, dpi=180)
        plt.close(fig)
        return output

    def build_plan_comparison_chart(
        self,
        monthly_rows: list[dict[str, Any]],
        filename: str,
        planned_metric: str,
        actual_metric: str,
        title: str,
    ) -> Path:
        return self.build_bar_comparison_chart(
            monthly_rows,
            filename,
            planned_metric,
            actual_metric,
            title,
            "Plan",
            "Actual",
            empty_message="No plan comparison data",
        )

    def build_bar_comparison_chart(
        self,
        monthly_rows: list[dict[str, Any]],
        filename: str,
        left_metric: str,
        right_metric: str,
        title: str,
        left_label: str,
        right_label: str,
        empty_message: str = "No comparison data",
    ) -> Path:
        output = self.charts_dir / filename
        if not monthly_rows:
            return self._plot_empty(output, empty_message)

        df = pd.DataFrame(monthly_rows)
        if df.empty or "month_label" not in df.columns:
            return self._plot_empty(output, empty_message)

        labels = df["month_label"].tolist()
        x_positions = range(len(labels))
        width = 0.36
        left_values = df[left_metric].fillna(0) if left_metric in df.columns else [0] * len(labels)
        right_values = df[right_metric].fillna(0) if right_metric in df.columns else [0] * len(labels)

        fig, ax = plt.subplots(figsize=(7.8, 3.9))
        ax.bar(
            [position - width / 2 for position in x_positions],
            left_values,
            width=width,
            color="#94A3B8",
            label=left_label,
        )
        ax.bar(
            [position + width / 2 for position in x_positions],
            right_values,
            width=width,
            color="#0C5460",
            label=right_label,
        )
        ax.set_title(title, fontsize=14)
        ax.set_xticks(list(x_positions), labels)
        ax.tick_params(axis="both", labelsize=9)
        ax.grid(axis="y", alpha=0.2)
        ax.legend(fontsize=9, loc="upper left")
        plt.tight_layout()
        fig.savefig(output, dpi=180)
        plt.close(fig)
        return output

    def _render_slide(self, slide_spec: dict[str, Any]) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg

        slide_type = slide_spec.get("type")
        if slide_type == "cover":
            self._add_title(slide, slide_spec.get("title", ""), top=1.8, size=30)
            self._add_subtitle(slide, slide_spec.get("subtitle", ""), top=2.65, size=18)
            self._add_center_badge(slide, slide_spec.get("client_name", "Wightlink"))
            return
        if slide_type == "divider":
            self._add_title(slide, slide_spec.get("title", ""), top=2.7, size=34)
            return
        if slide_type == "closing":
            self._add_title(slide, slide_spec.get("title", "Any Questions?"), top=2.8, size=34)
            return

        self._add_title(slide, slide_spec.get("title", ""))
        if slide_spec.get("subtitle"):
            self._add_subtitle(slide, slide_spec.get("subtitle", ""))

        charts = slide_spec.get("charts", [])
        table_rows = slide_spec.get("table", {}).get("rows", [])
        bullets = slide_spec.get("bullets", [])

        if slide_type == "agenda":
            self._add_bullets(slide, bullets, Inches(1.0), Inches(1.6), Inches(11.0), Inches(4.8), font_size=22)
        elif slide_type == "dual_chart_bullets":
            if len(charts) >= 1:
                slide.shapes.add_picture(str(charts[0]["path"]), Inches(0.45), Inches(1.35), width=Inches(6.0), height=Inches(3.15))
            if len(charts) >= 2:
                slide.shapes.add_picture(str(charts[1]["path"]), Inches(6.85), Inches(1.35), width=Inches(6.0), height=Inches(3.15))
            self._add_bullets(slide, bullets, Inches(0.8), Inches(4.85), Inches(11.8), Inches(1.65))
        elif slide_type == "kpi_cards_bullets":
            self._render_kpi_cards(slide, slide_spec.get("kpis", []), Inches(0.65), Inches(1.35), Inches(12.05), Inches(3.45))
            self._add_bullets(slide, bullets, Inches(0.8), Inches(5.15), Inches(11.8), Inches(1.25), font_size=14)
        elif slide_type == "single_chart_bullets":
            if charts:
                slide.shapes.add_picture(str(charts[0]["path"]), Inches(0.6), Inches(1.5), width=Inches(6.7), height=Inches(3.8))
            self._add_bullets(slide, bullets, Inches(7.65), Inches(1.6), Inches(4.6), Inches(3.7))
        elif slide_type == "wide_chart_bullets":
            if charts:
                slide.shapes.add_picture(str(charts[0]["path"]), Inches(1.0), Inches(1.28), width=Inches(11.3), height=Inches(4.55))
            self._add_bullets(slide, bullets, Inches(1.05), Inches(6.0), Inches(11.0), Inches(0.7), font_size=13)
        elif slide_type == "table_bullets":
            self._render_table(slide, table_rows, Inches(0.4), Inches(1.35), Inches(12.45), Inches(3.35))
            self._add_bullets(slide, bullets, Inches(0.8), Inches(4.95), Inches(11.8), Inches(1.45))
        elif slide_type == "dual_chart_table_bullets":
            if len(charts) >= 1:
                slide.shapes.add_picture(str(charts[0]["path"]), Inches(0.45), Inches(1.3), width=Inches(6.0), height=Inches(2.7))
            if len(charts) >= 2:
                slide.shapes.add_picture(str(charts[1]["path"]), Inches(6.85), Inches(1.3), width=Inches(6.0), height=Inches(2.7))
            self._render_table(slide, table_rows, Inches(0.45), Inches(4.2), Inches(8.0), Inches(2.15))
            self._add_bullets(slide, bullets, Inches(8.7), Inches(4.25), Inches(4.1), Inches(2.05), font_size=14)
        elif slide_type == "table_only":
            self._render_table(slide, table_rows, Inches(0.4), Inches(1.35), Inches(12.45), Inches(4.8))
            self._add_bullets(slide, bullets, Inches(0.8), Inches(6.15), Inches(11.8), Inches(0.7), font_size=13)
        elif slide_type == "bullets_only":
            self._add_bullets(slide, bullets, Inches(0.9), Inches(1.65), Inches(11.7), Inches(4.8), font_size=20)
        elif slide_type == "image_bullets":
            image_path = slide_spec.get("image_path")
            if image_path and Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), Inches(0.55), Inches(1.45), width=Inches(6.8), height=Inches(4.3))
            else:
                self._add_image_placeholder(slide, Inches(0.55), Inches(1.45), Inches(6.8), Inches(4.3), "Image / screenshot placeholder")
            self._add_bullets(slide, bullets, Inches(7.7), Inches(1.6), Inches(4.8), Inches(4.0))

        if slide_spec.get("source_note"):
            self._add_source_note(slide, slide_spec["source_note"])

    def _add_title(self, slide, text: str, top: float = 0.28, size: int = 24) -> None:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.0), Inches(0.6)).text_frame
        run = box.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = self.text_primary

    def _add_subtitle(self, slide, text: str, top: float = 0.88, size: int = 14) -> None:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.0), Inches(0.35)).text_frame
        run = box.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = self.text_secondary

    def _add_center_badge(self, slide, text: str) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(3.6), Inches(4.9), Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.light_fill
        shape.line.color.rgb = self.accent
        frame = shape.text_frame
        frame.clear()
        para = frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = self.accent

    def _add_bullets(self, slide, bullets: list[str], left, top, width, height, font_size: int = 16) -> None:
        text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
        text_frame.clear()
        bullet_list = bullets or ["No narrative was available for this slide."]
        for index, bullet in enumerate(bullet_list):
            para = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            para.text = bullet
            para.level = 0
            para.font.size = Pt(font_size)
            para.font.color.rgb = self.text_body
            para.space_after = Pt(10)

    def _render_table(self, slide, rows: list[dict[str, Any]], left, top, width, height) -> None:
        if not rows:
            self._add_image_placeholder(slide, left, top, width, height, "No table data")
            return
        headers = list(rows[0].keys())
        shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
        table = shape.table
        for col_index, header in enumerate(headers):
            cell = table.cell(0, col_index)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.light_fill
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(9)
            p.font.color.rgb = self.text_primary
        for row_index, row in enumerate(rows, start=1):
            for col_index, header in enumerate(headers):
                cell = table.cell(row_index, col_index)
                cell.text = str(row.get(header, ""))
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8.5)
                p.font.color.rgb = self.text_body

    def _render_kpi_cards(self, slide, kpis: list[dict[str, Any]], left, top, width, height) -> None:
        cards = list(kpis)
        if not cards:
            self._add_image_placeholder(slide, left, top, width, height, "No KPI data")
            return

        columns = 3
        rows = 2
        gap_x = Inches(0.2)
        gap_y = Inches(0.2)
        card_width = int((width - gap_x * (columns - 1)) / columns)
        card_height = int((height - gap_y * (rows - 1)) / rows)

        for index, kpi in enumerate(cards[: columns * rows]):
            row = index // columns
            column = index % columns
            card_left = left + column * (card_width + gap_x)
            card_top = top + row * (card_height + gap_y)

            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = self.light_fill
            card.line.width = Pt(1)

            accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, card_left, card_top, card_width, Inches(0.05))
            accent.fill.solid()
            accent.fill.fore_color.rgb = self.accent
            accent.line.color.rgb = self.accent

            label_frame = slide.shapes.add_textbox(card_left + Inches(0.16), card_top + Inches(0.14), card_width - Inches(0.32), Inches(0.28)).text_frame
            label_run = label_frame.paragraphs[0].add_run()
            label_run.text = str(kpi.get("label", "Metric"))
            label_run.font.size = Pt(10)
            label_run.font.bold = True
            label_run.font.color.rgb = self.text_secondary

            value_frame = slide.shapes.add_textbox(card_left + Inches(0.16), card_top + Inches(0.46), card_width - Inches(0.32), Inches(0.44)).text_frame
            value_run = value_frame.paragraphs[0].add_run()
            value_run.text = str(kpi.get("value", "--"))
            value_run.font.size = Pt(20)
            value_run.font.bold = True
            value_run.font.color.rgb = self.text_primary

            context_items = list(kpi.get("context_items", []))
            if not context_items:
                context_items = [{"text": line, "value": kpi.get("yoy")} for line in (list(kpi.get("context", [])) or [f"YoY: {kpi.get('yoy_label', '--')}"])]
            context_frame = slide.shapes.add_textbox(card_left + Inches(0.16), card_top + card_height - Inches(0.55), card_width - Inches(0.32), Inches(0.36)).text_frame
            context_frame.clear()
            for line_index, item in enumerate(context_items[:2]):
                para = context_frame.paragraphs[0] if line_index == 0 else context_frame.add_paragraph()
                para.text = str(item.get("text", ""))
                para.font.size = Pt(9)
                para.font.bold = True
                para.font.color.rgb = self._resolve_delta_color(str(kpi.get("key", "")), item.get("value"))

    def _resolve_delta_color(self, key: str, value: Any) -> RGBColor:
        if value is None or (isinstance(value, float) and pd.isna(value)):
            return self.text_secondary
        lower_is_better = key in {"cost", "cpa", "cpc"}
        positive = float(value) >= 0
        if positive != lower_is_better:
            return RGBColor(22, 101, 52)
        return RGBColor(185, 28, 28)

    def _add_image_placeholder(self, slide, left, top, width, height, text: str) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(236, 241, 246)
        shape.line.color.rgb = self.text_secondary
        frame = shape.text_frame
        frame.clear()
        para = frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(16)
        run.font.color.rgb = self.text_secondary

    def _add_source_note(self, slide, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.25)).text_frame
        run = box.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(8)
        run.font.color.rgb = self.text_secondary

    def _plot_empty(self, output: Path, message: str) -> Path:
        fig, ax = plt.subplots(figsize=(7.8, 4.0))
        ax.text(0.5, 0.5, message, ha="center", va="center", fontsize=16)
        ax.axis("off")
        plt.tight_layout()
        fig.savefig(output, dpi=180)
        plt.close(fig)
        return output


def _compact_currency_tick(value: float, _: Any) -> str:
    return _compact_currency_label(value)


def _compact_currency_label(value: float) -> str:
    absolute = abs(float(value))
    if absolute >= 1_000_000:
        return f"£{float(value) / 1_000_000:.1f}m"
    if absolute >= 1_000:
        return f"£{float(value) / 1_000:.0f}k"
    return f"£{float(value):,.0f}"
