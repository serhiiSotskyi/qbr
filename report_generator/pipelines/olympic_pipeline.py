from __future__ import annotations

import os
from pathlib import Path
from typing import Any, Iterable

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from report_generator.insights.olympic_insights import (
    analyze_atc,
    analyze_channels,
    analyze_trends,
    generate_executive_summary,
)
from src.auction_loader import load_auction_csv
from src.auction_metrics import format_auction_table, summarize_auction_insights
from src.trends_loader import TrendsLoader

os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib")
import matplotlib

matplotlib.use("Agg", force=True)
import matplotlib.pyplot as plt


def generate_olympic_report(
    rows,
    client_config,
    output_path,
    manual_inputs=None,
    trends_dir=None,
    auction_csv=None,
):
    config = dict(client_config or {})
    project_root = Path(config.get("_project_root", Path.cwd()))
    chart_styles = dict(config.get("_chart_styles", {}))
    template_path = _resolve_template_path(project_root, config)
    client_id = str(config.get("id", "olympic_holidays"))
    manual_content = _normalize_manual_inputs(manual_inputs)

    df = _coerce_rows_to_dataframe(rows)
    data = _prepare_datasets(df)
    insights = _build_performance_insights(data)

    trend_sections = _build_trend_sections(trends_dir, data, project_root, client_id, chart_styles)
    auction_section = _build_auction_section(auction_csv, config)

    requested_output = Path(output_path)
    if requested_output.suffix.lower() == ".txt":
        text_path = requested_output
        pptx_path = project_root / "output" / f"{client_id}_report.pptx"
    else:
        pptx_path = requested_output if requested_output.suffix.lower() == ".pptx" else requested_output.with_suffix(".pptx")
        text_path = project_root / "reports" / f"{client_id}_report.txt"

    charts_dir = project_root / "charts" / client_id / "qbr"
    charts_dir.mkdir(parents=True, exist_ok=True)
    chart_paths = _build_performance_charts(data, charts_dir, chart_styles)

    _build_presentation(
        pptx_path=pptx_path,
        template_path=template_path,
        data=data,
        insights=insights,
        trend_sections=trend_sections,
        auction_section=auction_section,
        manual_content=manual_content,
        chart_paths=chart_paths,
        client_name=str(config.get("name", "Olympic Holidays")),
        chart_styles=chart_styles,
    )

    text = _build_text_report(
        data=data,
        insights=insights,
        trend_sections=trend_sections,
        auction_section=auction_section,
        manual_content=manual_content,
    )
    text_path.parent.mkdir(parents=True, exist_ok=True)
    text_path.write_text(text, encoding="utf-8")

    return {
        "pptx_path": pptx_path,
        "text_path": text_path,
        "text": text,
        "data": data,
        "insights": insights,
        "trend_sections": trend_sections,
        "auction_section": auction_section,
    }


def _coerce_rows_to_dataframe(rows) -> pd.DataFrame:
    if isinstance(rows, pd.DataFrame):
        df = rows.copy()
    elif isinstance(rows, Iterable):
        df = pd.DataFrame(list(rows))
    else:
        raise ValueError("Olympic pipeline expects DataFrame-like input.")

    if df.empty:
        raise ValueError("Olympic pipeline received no rows.")
    return df


def _prepare_datasets(df: pd.DataFrame) -> dict[str, Any]:
    working_df = df.rename(
        columns={
            "Date": "date",
            "Campaign Type": "channel",
            "Purchases": "purchases",
            "Revenue": "revenue",
            "Cost": "cost",
            "Add to cart": "add_to_cart",
        }
    ).copy()

    required = {"date", "channel", "purchases", "revenue", "cost", "add_to_cart"}
    missing = required - set(working_df.columns)
    if missing:
        raise ValueError(f"Olympic Holidays CSV missing required columns: {sorted(missing)}")

    working_df = working_df[list(required)].copy()
    working_df["date"] = pd.to_datetime(working_df["date"], dayfirst=True, errors="coerce")
    working_df = working_df.dropna(subset=["date"]).copy()
    if working_df.empty:
        raise ValueError("Olympic Holidays CSV has no valid dates.")

    for column in ["purchases", "revenue", "cost", "add_to_cart"]:
        working_df[column] = pd.to_numeric(working_df[column], errors="coerce").fillna(0.0)
    working_df["channel"] = working_df["channel"].fillna("Unknown").astype(str).str.strip()
    working_df["month_start"] = working_df["date"].dt.to_period("M").dt.to_timestamp()
    working_df["month_label"] = working_df["month_start"].dt.strftime("%b %Y")
    working_df["year"] = working_df["date"].dt.year
    working_df["quarter"] = working_df["date"].dt.quarter

    all_monthly = (
        working_df.groupby("month_start", as_index=False)[["revenue", "cost", "purchases", "add_to_cart"]]
        .sum()
        .sort_values("month_start")
        .reset_index(drop=True)
    )
    all_monthly["month_label"] = all_monthly["month_start"].dt.strftime("%b %Y")
    all_monthly["cpa"] = _safe_divide(all_monthly["cost"], all_monthly["purchases"])
    all_monthly["aov"] = _safe_divide(all_monthly["revenue"], all_monthly["purchases"])
    all_monthly["cpatc"] = _safe_divide(all_monthly["cost"], all_monthly["add_to_cart"])
    all_monthly["year"] = all_monthly["month_start"].dt.year
    all_monthly["quarter"] = all_monthly["month_start"].dt.quarter

    selected_year, selected_quarter = _detect_latest_complete_quarter(all_monthly)
    quarter_monthly = all_monthly[(all_monthly["year"] == selected_year) & (all_monthly["quarter"] == selected_quarter)].copy()
    quarter_monthly = quarter_monthly.reset_index(drop=True)

    quarter_rows = working_df[
        (working_df["date"].dt.year == selected_year) & (working_df["date"].dt.quarter == selected_quarter)
    ].copy()

    channel_breakdown = (
        quarter_rows.groupby("channel", as_index=False)[["revenue", "cost", "purchases", "add_to_cart"]]
        .sum()
        .sort_values(["revenue", "cost"], ascending=False)
        .reset_index(drop=True)
    )
    channel_breakdown["cpa"] = _safe_divide(channel_breakdown["cost"], channel_breakdown["purchases"])
    channel_breakdown["aov"] = _safe_divide(channel_breakdown["revenue"], channel_breakdown["purchases"])
    channel_breakdown["cpatc"] = _safe_divide(channel_breakdown["cost"], channel_breakdown["add_to_cart"])
    channel_breakdown["cost_share"] = _share(channel_breakdown["cost"])
    channel_breakdown["revenue_share"] = _share(channel_breakdown["revenue"])

    atc_trends = quarter_monthly[["month_start", "month_label", "add_to_cart", "cpatc", "purchases"]].copy()
    atc_trends = atc_trends.rename(columns={"add_to_cart": "atc"})

    yearly = (
        working_df.groupby("year", as_index=False)[["revenue", "cost", "purchases", "add_to_cart"]]
        .sum()
        .sort_values("year")
        .reset_index(drop=True)
    )
    yearly["cpa"] = _safe_divide(yearly["cost"], yearly["purchases"])
    yearly["aov"] = _safe_divide(yearly["revenue"], yearly["purchases"])
    yearly["cpatc"] = _safe_divide(yearly["cost"], yearly["add_to_cart"])

    yoy = _build_quarter_yoy_summary(all_monthly, selected_year, selected_quarter)
    summary = {
        "selected_year": selected_year,
        "selected_quarter": selected_quarter,
        "period_start": quarter_monthly.iloc[0]["month_label"],
        "period_end": quarter_monthly.iloc[-1]["month_label"],
        "overall_revenue": float(quarter_monthly["revenue"].sum()),
        "overall_cost": float(quarter_monthly["cost"].sum()),
        "overall_purchases": float(quarter_monthly["purchases"].sum()),
        "overall_cpa": float(_safe_scalar(quarter_monthly["cost"].sum(), quarter_monthly["purchases"].sum())),
        "latest_revenue_change_pct": _pct_change(quarter_monthly.iloc[-2]["revenue"], quarter_monthly.iloc[-1]["revenue"]) if len(quarter_monthly) > 1 else 0.0,
        "latest_cost_change_pct": _pct_change(quarter_monthly.iloc[-2]["cost"], quarter_monthly.iloc[-1]["cost"]) if len(quarter_monthly) > 1 else 0.0,
    }

    correlations = {
        "revenue_cost": _correlation(quarter_monthly["revenue"], quarter_monthly["cost"]),
        "revenue_atc": _correlation(quarter_monthly["revenue"], quarter_monthly["add_to_cart"]),
    }

    return {
        "all_monthly": all_monthly,
        "monthly_performance": quarter_monthly,
        "channel_breakdown": channel_breakdown,
        "atc_trends": atc_trends,
        "end_period": quarter_monthly.tail(min(2, len(quarter_monthly))).reset_index(drop=True),
        "yearly_aggregates": yearly,
        "yoy_summary": yoy,
        "summary": summary,
        "correlations": correlations,
    }


def _build_performance_insights(data: dict[str, Any]) -> dict[str, Any]:
    executive = generate_executive_summary(data)
    trends = analyze_trends(data)
    channels = analyze_channels(data)
    atc = analyze_atc(data)
    return {
        "executive_summary": executive["bullets"],
        "overall_performance": trends["overall"],
        "end_of_period": trends["end_of_period"],
        "channel_performance": channels["bullets"],
        "atc_analysis": atc["bullets"],
    }


def _build_trend_sections(
    trends_dir: str | Path | None,
    data: dict[str, Any],
    project_root: Path,
    client_id: str,
    chart_styles: dict[str, Any],
) -> list[dict[str, Any]]:
    loader = TrendsLoader(trends_dir)
    trends_df = loader.load_from_directory()
    if trends_df.empty:
        return []

    charts_dir = project_root / "charts" / client_id / "market_trends"
    charts_dir.mkdir(parents=True, exist_ok=True)
    sections: list[dict[str, Any]] = []
    for source_file, source_df in trends_df.groupby("source_file"):
        monthly = (
            source_df.groupby(["month_start", "term"], as_index=False)["value"]
            .mean()
            .sort_values(["month_start", "term"])
            .reset_index(drop=True)
        )
        if monthly.empty:
            continue

        title = _title_from_filename(source_file)
        chart_path = _plot_market_trend_chart(charts_dir / f"{_slug(source_file)}.png", monthly, chart_styles)
        bullets = _describe_trend_source(monthly, data["summary"]["selected_quarter"])
        sections.append(
            {
                "title": title,
                "subtitle": "Google Trends",
                "chart_path": chart_path,
                "bullets": bullets,
            }
        )
    return sections


def _build_auction_section(auction_csv: str | Path | None, client_config: dict[str, Any]) -> dict[str, Any] | None:
    if not auction_csv:
        return None
    try:
        auction_df = load_auction_csv(auction_csv)
    except Exception:
        return None
    if auction_df.empty:
        return None

    summary = summarize_auction_insights(
        auction_df,
        client_domain=client_config.get("site_domain") or client_config.get("auction_insights", {}).get("client_domain"),
        known_competitors=client_config.get("auction_insights", {}).get("known_competitors", []),
    )
    if not summary:
        return None

    bullets = []
    if summary.get("our_impression_share") is not None:
        bullets.append(f"Our impression share in the uploaded auction report was {summary['our_impression_share'] * 100:.1f}%.")
    if summary.get("top_impression_share_competitors"):
        top = summary["top_impression_share_competitors"][0]
        bullets.append(f"{top['domain']} had the highest impression share at {top['value'] * 100:.1f}%.")
    if summary.get("top_overlap_competitors"):
        top = summary["top_overlap_competitors"][0]
        bullets.append(f"{top['domain']} had the highest overlap rate at {top['value'] * 100:.1f}%.")
    if summary.get("average_top_of_page_rate") is not None:
        bullets.append(f"The average top of page rate across competitors was {summary['average_top_of_page_rate'] * 100:.1f}%.")

    return {
        "title": "Auction Insights",
        "subtitle": "Uploaded auction report",
        "table": format_auction_table(auction_df).head(8),
        "bullets": bullets[:4],
    }


def _build_text_report(
    data: dict[str, Any],
    insights: dict[str, Any],
    trend_sections: list[dict[str, Any]],
    auction_section: dict[str, Any] | None,
    manual_content: dict[str, list[str]],
) -> str:
    sections: list[tuple[str, list[str]]] = [
        ("Executive Summary", insights["executive_summary"]),
    ]
    if trend_sections:
        market_trend_bullets = []
        for item in trend_sections:
            market_trend_bullets.append(f"{item['title']}: {item['bullets'][0]}")
            if len(item["bullets"]) > 1:
                market_trend_bullets.append(item["bullets"][1])
        sections.append(("Market Trends", market_trend_bullets[:6]))
    if auction_section:
        sections.append(("Auction Insights", auction_section["bullets"]))
    sections.extend(
        [
            ("Overall Performance", insights["overall_performance"]),
            ("End of Period", insights["end_of_period"]),
            ("Channel Performance", insights["channel_performance"]),
            ("ATC Analysis", insights["atc_analysis"]),
        ]
    )
    if manual_content["actions"]:
        sections.append(("Actions", manual_content["actions"]))
    if manual_content["opportunities"]:
        sections.append(("Opportunities", manual_content["opportunities"]))

    blocks = []
    for title, bullets in sections:
        clean_bullets = [str(bullet).strip() for bullet in bullets if str(bullet).strip()]
        blocks.append("\n".join([f"[{title}]"] + [f"- {bullet}" for bullet in clean_bullets]))
    return "\n\n".join(blocks).strip() + "\n"


def _build_presentation(
    pptx_path: Path,
    template_path: Path,
    data: dict[str, Any],
    insights: dict[str, Any],
    trend_sections: list[dict[str, Any]],
    auction_section: dict[str, Any] | None,
    manual_content: dict[str, list[str]],
    chart_paths: dict[str, Path],
    client_name: str,
    chart_styles: dict[str, Any],
) -> None:
    template_prs = Presentation(str(template_path))
    prs = Presentation()
    prs.slide_width = template_prs.slide_width
    prs.slide_height = template_prs.slide_height

    colors = (chart_styles or {}).get("colors", {})
    palette = {
        "title": _hex_to_rgb(colors.get("text_primary", "#142A4D")),
        "subtitle": _hex_to_rgb(colors.get("text_secondary", "#5A5A5A")),
        "body": _hex_to_rgb(colors.get("text_body", "#2D2D2D")),
        "header": _hex_to_rgb(colors.get("table_header", "#DBE6F4")),
        "total": _hex_to_rgb(colors.get("table_total", "#EBF1FA")),
    }

    quarter_label = f"Q{data['summary']['selected_quarter']} {data['summary']['selected_year']}"
    subtitle = f"{client_name} | {quarter_label}"

    _add_text_slide(prs, "Executive Summary", subtitle, insights["executive_summary"], palette)
    for trend_section in trend_sections:
        _add_chart_slide(
            prs,
            title=trend_section["title"],
            subtitle=trend_section["subtitle"],
            chart_path=trend_section["chart_path"],
            bullets=trend_section["bullets"],
            palette=palette,
        )
    if auction_section:
        _add_table_slide(
            prs,
            title=auction_section["title"],
            subtitle=auction_section["subtitle"],
            table_df=auction_section["table"],
            bullets=auction_section["bullets"],
            palette=palette,
        )
    _add_chart_slide(prs, "Overall Performance", subtitle, chart_paths["overall_performance"], insights["overall_performance"], palette)
    _add_chart_slide(prs, "End of Period", "Latest months", chart_paths["end_of_period"], insights["end_of_period"], palette)
    _add_dual_chart_slide(
        prs,
        "Channel Performance",
        "Cost and revenue share",
        chart_paths["channel_cost_share"],
        chart_paths["channel_revenue_share"],
        insights["channel_performance"],
        palette,
    )
    _add_chart_slide(prs, "ATC Analysis", "Add to cart and CPATC", chart_paths["atc_analysis"], insights["atc_analysis"], palette)
    if manual_content["actions"]:
        _add_text_slide(prs, "Actions", "Manual input", manual_content["actions"], palette)
    if manual_content["opportunities"]:
        _add_text_slide(prs, "Opportunities", "Manual input", manual_content["opportunities"], palette)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))


def _build_performance_charts(data: dict[str, Any], charts_dir: Path, chart_styles: dict[str, Any]) -> dict[str, Path]:
    colors = (chart_styles or {}).get("colors", {})
    revenue_color = colors.get("cost", "#14B8A6")
    cpa_color = colors.get("cvr", "#1D4ED8")
    atc_color = colors.get("leads", "#0F172A")
    return {
        "overall_performance": _plot_bar_line(
            charts_dir / "overall_performance.png",
            data["monthly_performance"],
            bar_col="revenue",
            line_col="cpa",
            bar_label="Revenue",
            line_label="CPA",
            bar_color=revenue_color,
            line_color=cpa_color,
        ),
        "end_of_period": _plot_bar_line(
            charts_dir / "end_of_period.png",
            data["end_period"],
            bar_col="revenue",
            line_col="cpa",
            bar_label="Revenue",
            line_label="CPA",
            bar_color="#F97316",
            line_color=cpa_color,
        ),
        "channel_cost_share": _plot_pie(
            charts_dir / "channel_cost_share.png",
            data["channel_breakdown"],
            value_col="cost",
            label_col="channel",
            colors=["#0F766E", "#14B8A6", "#2DD4BF", "#99F6E4", "#CCFBF1"],
        ),
        "channel_revenue_share": _plot_pie(
            charts_dir / "channel_revenue_share.png",
            data["channel_breakdown"],
            value_col="revenue",
            label_col="channel",
            colors=["#1D4ED8", "#2563EB", "#60A5FA", "#93C5FD", "#DBEAFE"],
        ),
        "atc_analysis": _plot_bar_line(
            charts_dir / "atc_analysis.png",
            data["atc_trends"],
            bar_col="atc",
            line_col="cpatc",
            bar_label="Add to Cart",
            line_label="CPATC",
            bar_color=atc_color,
            line_color=cpa_color,
        ),
    }


def _add_chart_slide(prs, title: str, subtitle: str, chart_path: Path, bullets: list[str], palette: dict[str, RGBColor]) -> None:
    slide = _new_slide(prs)
    _add_title(slide, title, palette["title"])
    _add_subtitle(slide, subtitle, palette["subtitle"])
    slide.shapes.add_picture(str(chart_path), Inches(0.6), Inches(1.3), width=Inches(7.2), height=Inches(4.8))
    _add_bullets(slide, bullets, left=8.15, top=1.45, width=4.2, height=4.8, text_color=palette["body"])


def _add_text_slide(prs, title: str, subtitle: str, bullets: list[str], palette: dict[str, RGBColor]) -> None:
    slide = _new_slide(prs)
    _add_title(slide, title, palette["title"])
    _add_subtitle(slide, subtitle, palette["subtitle"])
    _add_bullets(slide, bullets, left=0.9, top=1.5, width=11.2, height=4.8, text_color=palette["body"])


def _add_table_slide(
    prs,
    title: str,
    subtitle: str,
    table_df: pd.DataFrame,
    bullets: list[str],
    palette: dict[str, RGBColor],
) -> None:
    slide = _new_slide(prs)
    _add_title(slide, title, palette["title"])
    _add_subtitle(slide, subtitle, palette["subtitle"])
    _render_table(slide, table_df, Inches(0.45), Inches(1.35), Inches(8.2), Inches(3.7), palette)
    _add_bullets(slide, bullets, left=8.95, top=1.45, width=3.25, height=4.6, text_color=palette["body"])


def _add_dual_chart_slide(prs, title: str, subtitle: str, left_chart: Path, right_chart: Path, bullets: list[str], palette: dict[str, RGBColor]) -> None:
    slide = _new_slide(prs)
    _add_title(slide, title, palette["title"])
    _add_subtitle(slide, subtitle, palette["subtitle"])
    slide.shapes.add_picture(str(left_chart), Inches(0.45), Inches(1.35), width=Inches(4.2), height=Inches(3.8))
    slide.shapes.add_picture(str(right_chart), Inches(4.85), Inches(1.35), width=Inches(4.2), height=Inches(3.8))
    _add_bullets(slide, bullets, left=9.3, top=1.45, width=2.95, height=4.6, text_color=palette["body"])


def _render_table(slide, table_df: pd.DataFrame, left, top, width, height, palette: dict[str, RGBColor]) -> None:
    safe_df = table_df.copy()
    if safe_df.empty:
        safe_df = pd.DataFrame([{"Status": "No data"}])
    table_shape = slide.shapes.add_table(len(safe_df) + 1, len(safe_df.columns), left, top, width, height)
    table = table_shape.table
    for col_idx, col_name in enumerate(safe_df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        _style_cell(cell, True, palette["title"])
        cell.fill.solid()
        cell.fill.fore_color.rgb = palette["header"]
    for row_idx, (_, row) in enumerate(safe_df.iterrows(), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            _style_cell(cell, False, palette["title"])


def _style_cell(cell, bold: bool, color: RGBColor) -> None:
    para = cell.text_frame.paragraphs[0]
    if not para.runs:
        para.add_run()
    para.runs[0].font.bold = bold
    para.runs[0].font.size = Pt(10)
    para.runs[0].font.color.rgb = color


def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(247, 249, 252)
    return slide


def _add_title(slide, text: str, color: RGBColor) -> None:
    text_frame = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(12.0), Inches(0.6)).text_frame
    run = text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = color


def _add_subtitle(slide, text: str, color: RGBColor) -> None:
    text_frame = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12.0), Inches(0.3)).text_frame
    run = text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.color.rgb = color


def _add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float, text_color: RGBColor) -> None:
    text_frame = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    for idx, bullet in enumerate((bullets or ["No data available"])[:6]):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = str(bullet)
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = text_color


def _plot_bar_line(output_path: Path, frame: pd.DataFrame, bar_col: str, line_col: str, bar_label: str, line_label: str, bar_color: str, line_color: str) -> Path:
    fig, ax1 = plt.subplots(figsize=(8.0, 4.6))
    ax2 = ax1.twinx()
    ax1.bar(frame["month_label"], frame[bar_col], color=bar_color, alpha=0.85, label=bar_label)
    ax2.plot(frame["month_label"], frame[line_col], marker="o", linewidth=2.4, color=line_color, label=line_label)
    ax1.set_ylabel(bar_label)
    ax2.set_ylabel(line_label)
    ax1.tick_params(axis="x", rotation=35)
    ax1.grid(axis="y", alpha=0.2)
    lines_1, labels_1 = ax1.get_legend_handles_labels()
    lines_2, labels_2 = ax2.get_legend_handles_labels()
    ax1.legend(lines_1 + lines_2, labels_1 + labels_2, loc="upper left")
    plt.tight_layout()
    fig.savefig(output_path, dpi=180)
    plt.close(fig)
    return output_path


def _plot_pie(output_path: Path, frame: pd.DataFrame, value_col: str, label_col: str, colors: list[str]) -> Path:
    chart_df = frame[frame[value_col] > 0].copy()
    fig, ax = plt.subplots(figsize=(4.6, 3.9))
    if chart_df.empty:
        ax.text(0.5, 0.5, "No data", ha="center", va="center")
        ax.axis("off")
    else:
        ax.pie(
            chart_df[value_col],
            labels=chart_df[label_col],
            autopct=lambda pct: f"{pct:.1f}%" if pct >= 4 else "",
            startangle=90,
            colors=colors[: len(chart_df)],
            textprops={"fontsize": 9},
        )
        ax.axis("equal")
    plt.tight_layout()
    fig.savefig(output_path, dpi=180)
    plt.close(fig)
    return output_path


def _plot_market_trend_chart(output_path: Path, monthly: pd.DataFrame, chart_styles: dict[str, Any]) -> Path:
    colors = (chart_styles or {}).get("colors", {})
    palette = [
        colors.get("trend_current", "#0E7490"),
        colors.get("trend_prior", "#94A3B8"),
        "#1D4ED8",
        "#F97316",
        "#0F172A",
    ]
    fig, ax = plt.subplots(figsize=(8.0, 4.6))
    for idx, (term, term_df) in enumerate(monthly.groupby("term")):
        plot_df = term_df.sort_values("month_start")
        ax.plot(
            plot_df["month_start"].dt.strftime("%b %Y"),
            plot_df["value"],
            marker="o",
            linewidth=2.2,
            label=term,
            color=palette[idx % len(palette)],
        )
    ax.tick_params(axis="x", rotation=35)
    ax.set_ylabel("Interest")
    ax.grid(axis="y", alpha=0.2)
    ax.legend(loc="upper left", fontsize=9)
    plt.tight_layout()
    fig.savefig(output_path, dpi=180)
    plt.close(fig)
    return output_path


def _describe_trend_source(monthly: pd.DataFrame, selected_quarter: int) -> list[str]:
    total_monthly = monthly.groupby("month_start", as_index=False)["value"].sum().sort_values("month_start")
    total_monthly["year"] = total_monthly["month_start"].dt.year
    total_monthly["quarter"] = total_monthly["month_start"].dt.quarter
    bullets: list[str] = []

    years = sorted(total_monthly["year"].unique())
    if len(years) >= 2:
        current_year = years[-1]
        prior_year = years[-2]
        current = total_monthly[total_monthly["year"] == current_year].copy()
        prior = total_monthly[total_monthly["year"] == prior_year].copy()
        current["month_num"] = current["month_start"].dt.month
        prior["month_num"] = prior["month_start"].dt.month
        matched = current.merge(prior, on="month_num", suffixes=("_current", "_prior"))
        if not matched.empty:
            yoy = _pct_change(matched["value_prior"].sum(), matched["value_current"].sum())
            bullets.append(f"Search interest across this upload changed {yoy:.1f}% YoY across matched months.")

        q_current = total_monthly[(total_monthly["year"] == current_year) & (total_monthly["quarter"] == selected_quarter)]
        q_prior = total_monthly[(total_monthly["year"] == prior_year) & (total_monthly["quarter"] == selected_quarter)]
        if not q_current.empty and not q_prior.empty:
            q_yoy = _pct_change(q_prior["value"].sum(), q_current["value"].sum())
            bullets.append(f"Quarter-level search interest for Q{selected_quarter} changed {q_yoy:.1f}% versus the prior year.")

    peak = total_monthly.loc[total_monthly["value"].idxmax()]
    bullets.append(f"Peak interest occurred in {peak['month_start'].strftime('%b %Y')} with an indexed value of {peak['value']:.1f}.")
    latest = total_monthly.iloc[-1]
    previous = total_monthly.iloc[-2] if len(total_monthly) > 1 else latest
    bullets.append(
        f"Latest month interest {'rose' if latest['value'] >= previous['value'] else 'fell'} from {previous['value']:.1f} to {latest['value']:.1f}."
    )
    return bullets[:4]


def _detect_latest_complete_quarter(monthly: pd.DataFrame) -> tuple[int, int]:
    quarterly_months = (
        monthly.groupby(["year", "quarter"])["month_start"]
        .nunique()
        .reset_index(name="month_count")
        .sort_values(["year", "quarter"])
    )
    complete = quarterly_months[quarterly_months["month_count"] == 3]
    if complete.empty:
        latest = quarterly_months.iloc[-1]
        return int(latest["year"]), int(latest["quarter"])
    latest = complete.iloc[-1]
    return int(latest["year"]), int(latest["quarter"])


def _build_quarter_yoy_summary(all_monthly: pd.DataFrame, selected_year: int, selected_quarter: int) -> dict[str, Any] | None:
    current = all_monthly[(all_monthly["year"] == selected_year) & (all_monthly["quarter"] == selected_quarter)].copy()
    prior = all_monthly[(all_monthly["year"] == selected_year - 1) & (all_monthly["quarter"] == selected_quarter)].copy()
    if current.empty or prior.empty:
        return None
    current["month_num"] = current["month_start"].dt.month
    prior["month_num"] = prior["month_start"].dt.month
    matched = current.merge(prior, on="month_num", suffixes=("_current", "_prior"))
    if matched.empty:
        return None

    revenue_current = matched["revenue_current"].sum()
    revenue_prior = matched["revenue_prior"].sum()
    cost_current = matched["cost_current"].sum()
    cost_prior = matched["cost_prior"].sum()
    purchases_current = matched["purchases_current"].sum()
    purchases_prior = matched["purchases_prior"].sum()
    cpa_current = _safe_scalar(cost_current, purchases_current)
    cpa_prior = _safe_scalar(cost_prior, purchases_prior)
    return {
        "current_year": selected_year,
        "prior_year": selected_year - 1,
        "revenue_change_pct": _pct_change(revenue_prior, revenue_current),
        "cost_change_pct": _pct_change(cost_prior, cost_current),
        "purchases_change_pct": _pct_change(purchases_prior, purchases_current),
        "cpa_change_pct": _pct_change(cpa_prior, cpa_current),
    }


def _normalize_manual_inputs(manual_inputs) -> dict[str, list[str]]:
    payload = manual_inputs if isinstance(manual_inputs, dict) else {}
    normalized = {"actions": [], "opportunities": []}
    for key in normalized:
        value = payload.get(key, [])
        if isinstance(value, list):
            normalized[key] = [str(item).strip() for item in value if str(item).strip()]
    return normalized


def _resolve_template_path(project_root: Path, client_config: dict[str, Any]) -> Path:
    template_path = Path(str(client_config.get("template_path", "templates/report_template.pptx")))
    return template_path if template_path.is_absolute() else project_root / template_path


def _safe_divide(numerator, denominator):
    numerator_series = numerator.astype(float) if isinstance(numerator, pd.Series) else pd.Series(numerator, dtype=float)
    denominator_series = denominator.astype(float) if isinstance(denominator, pd.Series) else pd.Series(denominator, dtype=float)
    denominator_series = denominator_series.replace(0, float("nan"))
    return numerator_series.divide(denominator_series)


def _safe_scalar(numerator: float, denominator: float) -> float:
    if denominator == 0:
        return 0.0
    return float(numerator / denominator)


def _share(series: pd.Series) -> pd.Series:
    total = float(series.sum())
    if total == 0:
        return pd.Series([0.0] * len(series), index=series.index)
    return (series / total) * 100


def _pct_change(base: float, current: float) -> float:
    if base == 0 or pd.isna(base) or pd.isna(current):
        return 0.0
    return ((current - base) / base) * 100


def _correlation(left: pd.Series, right: pd.Series) -> float | None:
    if len(left) < 2 or len(right) < 2:
        return None
    value = left.corr(right)
    return None if pd.isna(value) else float(value)


def _title_from_filename(filename: str) -> str:
    stem = Path(filename).stem
    cleaned = stem.replace("_", " ").replace("-", " ").strip()
    return " ".join(part.capitalize() if not part.isupper() else part for part in cleaned.split())


def _slug(value: str) -> str:
    return value.lower().replace(" ", "_").replace("/", "_").replace("-", "_").replace(".", "_")


def _hex_to_rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    if len(value) != 6:
        value = "000000"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
