from __future__ import annotations

import csv
import json
import unittest
from datetime import datetime, timezone
from io import BytesIO, StringIO
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZipFile

from pptx import Presentation

from claude_handoff import DEFAULT_OLYMPIC_REFERENCE_PPTX_PATH, build_olympic_holidays_claude_handoff_package


GENERATED_AT = datetime(2026, 6, 30, 12, 0, tzinfo=timezone.utc)

SYNTHETIC_REPORT_TEXT = """[Cover]
Slide: 1
Subtitle: Q2 2026 (Apr - Jun 2026)
- Olympic Holidays Q2 2026 performance review.

[Olympic Holidays Trends]
Slide: 2
- Peak interest occurred in Jun 2026 with an indexed value of 62.0.

[Holidays to Greece Trends]
Slide: 3
- Peak interest occurred in May 2026 with an indexed value of 71.0.

[Executive Summary]
Slide: 4
- Q2 2026 delivered £512,345 revenue from 321.5 purchases on £88,765 spend.
- Quarter efficiency closed at £276 CPA, £24 cost per ATC, and £1,593 AOV.
- Versus Q2 2025, revenue moved +12.3%, cost moved +4.5%, and purchases moved +8.1%.

[Auction Insights]
Slide: 5
- The uploaded auction report included 3 competitor row(s).
Table:
Domain | Impression Share
-------+-----------------
you    | 51.0%

[Overall Performance]
Slide: 6
- Revenue grew across the quarter.

[YoY Matched-Month Comparison]
Slide: 7
- The comparison covers 3 of 3 month(s) in the quarter.
Table:
Metric | Q2 2025 | Q2 2026 | Change
-------+---------+---------+-------
Revenue | £456,000 | £512,345 | +12.3%

[End of Period]
Slide: 8
- Jun 2026 revenue was £201,000 versus £178,000 in May 2026.

[Channel Performance]
Slide: 9
- Brand drove the largest revenue share.
Table:
Channel | Revenue
--------+--------
Brand   | £300,000

[Generic Performance]
Slide: 10
- Generic delivered £100,000 revenue from £40,000 spend at £400 CPA.
Table:
Revenue | Cost | Purchases | Month
--------+------+-----------+------
£33,000 | £13,000 | 82.5 | Apr 2026

[ATC Analysis]
Slide: 11
- The quarter delivered 3654.1 total add-to-cart actions from £88,765 spend.
- Purchases totalled 321.5, providing the downstream conversion context for ATC performance.
Table:
Month | Add to Cart | Cost per ATC | Purchases
------+-------------+--------------+----------
Apr 2026 | 1100.0 | £25 | 100.0
"""


class OlympicClaudeHandoffPackageTests(unittest.TestCase):
    def test_builds_olympic_holidays_handoff_package_from_synthetic_inputs(self) -> None:
        package, manifest = _build_sample_handoff()

        expected_names = {
            "report.txt",
            "original_streamlit_prompt.txt",
            "olympic_holidays_streamlit_output.pptx",
            "reference_deck_exported_from_google_slides.pptx",
            "README_FOR_CLAUDE.txt",
            "CLAUDE_PROMPT.txt",
            "SLIDE_MAPPING.csv",
            "SOURCE_SECTION_INDEX.txt",
            "INPUT_FILES_MANIFEST.txt",
            "REFERENCE_DECK_OUTLINE.txt",
            "QA_CHECKLIST.txt",
            "CHART_QA_ADDENDUM_FOR_CLAUDE.txt",
            "PACKAGE_MANIFEST.json",
            "source_data/performance.csv",
            "source_data/auction_insights.csv",
            "source_data/google_trends_olympic_holidays.csv",
            "source_data/google_trends_category.csv",
        }
        self.assertTrue(expected_names.issubset(set(package.namelist())))
        self.assertNotIn("prompt.txt", package.namelist())

        self.assertEqual(manifest["report_family"], "Olympic Holidays PPC QBR")
        self.assertEqual(manifest["period_label"], "Q2 2026 (Apr - Jun 2026)")
        self.assertEqual(manifest["quarter_short"], "Q2 2026")
        self.assertEqual(manifest["period_months_label"], "April - June 2026")
        self.assertEqual(manifest["reference_slide_count"], 11)
        self.assertEqual(manifest["streamlit_slide_count"], 11)
        self.assertTrue(manifest["has_reference_pptx"])
        self.assertEqual(
            manifest["headline_kpis"],
            {
                "revenue": "£512,345",
                "purchases": "321.5",
                "cpa": "£276",
                "cost": "£88,765",
                "aov": "£1,593",
                "cost_per_atc": "£24",
                "total_atc": "3654.1",
            },
        )
        self.assertEqual(
            [(source["standard_name"], source["display_name"]) for source in manifest["trend_sources"]],
            [
                ("source_data/google_trends_olympic_holidays.csv", "Olympic Holidays"),
                ("source_data/google_trends_category.csv", "Holidays to Greece"),
            ],
        )

    def test_generated_guidance_references_chart_qa_and_uses_11_slide_mapping(self) -> None:
        package, _manifest = _build_sample_handoff()

        for filename in ("README_FOR_CLAUDE.txt", "CLAUDE_PROMPT.txt", "QA_CHECKLIST.txt"):
            content = package.read(filename).decode("utf-8")
            self.assertIn("CHART_QA_ADDENDUM_FOR_CLAUDE.txt", content)
            self.assertIn("Q2 2026 (Apr - Jun 2026)", content)

        slide_mapping = package.read("SLIDE_MAPPING.csv").decode("utf-8")
        rows = list(csv.DictReader(StringIO(slide_mapping)))
        self.assertEqual(len(rows), 11)
        self.assertIn("Olympic Holidays", slide_mapping)
        self.assertIn("Holidays to Greece", slide_mapping)
        self.assertIn("Revenue £512,345; Purchases 321.5; CPA £276", slide_mapping)

    def test_source_index_and_manifest_capture_source_metadata(self) -> None:
        package, manifest = _build_sample_handoff()

        source_index = package.read("SOURCE_SECTION_INDEX.txt").decode("utf-8")
        self.assertIn("Detected source sections: 11", source_index)
        self.assertIn("01. line 1, slide 1: [Cover]", source_index)
        self.assertIn("source_data/google_trends_olympic_holidays.csv: Olympic Holidays", source_index)
        self.assertIn("source_data/google_trends_category.csv: Holidays to Greece", source_index)

        input_manifest = package.read("INPUT_FILES_MANIFEST.txt").decode("utf-8")
        self.assertIn("Original upload: performance_export.csv", input_manifest)
        self.assertIn("Original upload: auction_insights.csv", input_manifest)
        self.assertIn("Display term: Olympic Holidays", input_manifest)
        self.assertIn("Display term: Holidays to Greece", input_manifest)

        manifest_from_zip = json.loads(package.read("PACKAGE_MANIFEST.json").decode("utf-8"))
        self.assertEqual(manifest_from_zip["generated_at"], "2026-06-30T12:00:00Z")
        self.assertEqual(manifest_from_zip["headline_kpis"], manifest["headline_kpis"])
        self.assertEqual(manifest_from_zip["placeholder_sections"], [])
        self.assertFalse(any("Q1 2026 delivered" in warning for warning in manifest_from_zip["warnings"]))


def _build_sample_handoff() -> tuple[ZipFile, dict]:
    with TemporaryDirectory() as tmpdir:
        root = Path(tmpdir)
        performance = root / "performance_export.csv"
        performance.write_text(
            "Date,Campaign Type,Revenue,Purchases,Cost,Add to cart,CPA,Cost per ATC,AOV\n"
            "01/04/2026,Brand,300000,200,40000,2200,200,18,1500\n",
            encoding="utf-8",
        )
        auction = root / "auction_insights.csv"
        auction.write_text("Domain,Impression share\nyou,51%\n", encoding="utf-8")
        brand_trend = root / "brand_export.csv"
        brand_trend.write_text('"Time","Olympic Holidays"\n"2026-04",54\n', encoding="utf-8")
        category_trend = root / "category_export.csv"
        category_trend.write_text('"Time","Holidays to Greece"\n"2026-04",66\n', encoding="utf-8")

        handoff_bytes, manifest = build_olympic_holidays_claude_handoff_package(
            report_text=SYNTHETIC_REPORT_TEXT,
            prompt_text="Original Streamlit prompt for Olympic Holidays.",
            generated_pptx=_make_pptx_bytes(slide_count=11),
            performance_csv=performance,
            auction_csv=auction,
            trend_csv_files=[category_trend, brand_trend],
            reference_pptx=DEFAULT_OLYMPIC_REFERENCE_PPTX_PATH,
            generated_at=GENERATED_AT,
        )

    return ZipFile(BytesIO(handoff_bytes)), manifest


def _make_pptx_bytes(*, slide_count: int) -> bytes:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    for _index in range(slide_count):
        presentation.slides.add_slide(blank_layout)
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()
