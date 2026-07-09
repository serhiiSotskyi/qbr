from __future__ import annotations

import csv
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from src.other_campaigns import load_other_campaign_summary
from src.report_pipeline import ReportPipeline
from utils.text_report import TextReportPipeline


ROOT = Path(__file__).resolve().parent.parent


class WendyWuOtherCampaignTests(unittest.TestCase):
    def test_other_campaign_parser_combines_ms_and_google_exports(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            _write_ms_campaign_export(root / "ms_campaigns.csv")
            _write_google_campaign_export(root / "google_campaigns.csv")
            _write_google_time_series_export(root / "google_daily.csv")

            summary = load_other_campaign_summary(
                root,
                exclude_terms=[
                    "brand",
                    "japan",
                    "china",
                    "india",
                    "se asia",
                    "vietnam",
                    "cambodia",
                    "thailand",
                    "malaysia",
                    "borneo",
                    "central asia",
                    "mongolia",
                ],
            )

        self.assertIsNotNone(summary)
        assert summary is not None
        self.assertEqual(summary["source_files"], ["google_campaigns.csv", "google_daily.csv", "ms_campaigns.csv"])
        self.assertFalse(summary["daily_impressions"].empty)

        top_click_names = summary["top_clicks"]["Campaign"].tolist()
        top_conversion_names = summary["top_conversions"]["Campaign"].tolist()
        self.assertNotIn("Mongolia", top_click_names)
        self.assertNotIn("Central Asia", top_conversion_names)
        self.assertIn("South Korea", top_conversion_names)

        other_raw_campaigns = " ".join(summary["other_campaign_rows"]["Raw Campaign"].tolist()).lower()
        self.assertNotIn("brand", other_raw_campaigns)
        self.assertNotIn("japan", other_raw_campaigns)
        self.assertNotIn("vietnam", other_raw_campaigns)
        self.assertNotIn("china", other_raw_campaigns)
        self.assertNotIn("central asia", other_raw_campaigns)
        self.assertNotIn("mongolia", other_raw_campaigns)

    def test_wendy_wu_pipeline_adds_other_top_campaigns_slide_from_uploads(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            performance_csv = _write_performance_csv(root / "performance.csv")
            other_dir = root / "other_campaigns"
            other_dir.mkdir()
            _write_ms_campaign_export(other_dir / "ms_campaigns.csv")
            output_pptx = root / "wendy_wu.pptx"

            pipeline = ReportPipeline(project_root=ROOT)
            pipeline.charts_root = root / "charts"
            generated_pptx = pipeline.run(
                input_csv=performance_csv,
                output_pptx=output_pptx,
                client_id="wendy_wu",
                other_campaigns_dir=other_dir,
            )

            text_pipeline = TextReportPipeline(project_root=ROOT)
            report_txt = root / "report.txt"
            text_pipeline.run(
                input_csv=performance_csv,
                output_txt=report_txt,
                client_id="wendy_wu",
                other_campaigns_dir=other_dir,
            )

            pptx_text = _pptx_text(generated_pptx)
            report_text = report_txt.read_text(encoding="utf-8")

        self.assertIn("Other (Destination) Top 10 campaigns", pptx_text)
        self.assertIn("Other (Destination) Top 10 campaigns", report_text)
        self.assertIn("Top 10 by Clicks", report_text)
        self.assertIn("South Korea", report_text)
        self.assertNotIn("Mongolia", report_text)
        self.assertNotIn("Central Asia", report_text)
        self.assertNotIn("UK - Brand - Core", report_text)
        self.assertNotIn("UK - Generic - Vietnam", report_text)

    def test_daily_time_series_without_campaign_names_does_not_create_summary(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            _write_google_time_series_export(root / "google_daily.csv")

            summary = load_other_campaign_summary(
                root,
                exclude_terms=[
                    "brand",
                    "japan",
                    "china",
                    "india",
                    "se asia",
                    "vietnam",
                    "cambodia",
                    "thailand",
                    "malaysia",
                    "borneo",
                    "central asia",
                    "mongolia",
                ],
            )

        self.assertIsNone(summary)


def _write_ms_campaign_export(path: Path) -> None:
    rows = [
        ["Campaign report (April 01, 2026 - June 30, 2026)"],
        ["Request Id: test"],
        [],
        ["Campaign ID", "Status", "Campaign", "Campaign Type", "Labels", "Clicks", "Impr.", "Spend", "Conv."],
        ["1", "Enabled", "UK - Brand - Core - Core", "Search", "Brand", "100", "1000", "10.00", "20"],
        ["2", "Enabled", "UK - Generic - Japan - General", "Search", "Asia; Japan", "90", "900", "90.00", "9"],
        ["3", "Enabled", "UK - Generic - Vietnam - General", "Search", "SE Asia; Vietnam", "80", "800", "80.00", "8"],
        ["4", "Enabled", "UK - Generic - Mongolia", "Search", "Asia; Other", "300", "3000", "120.00", "12"],
        ["5", "Enabled", "UK - Generic - Central Asia - General", "Search", "Central Asia; Other", "220", "2200", "110.00", "30"],
        ["6", "Enabled", "UK - Generic - South Korea - General", "Search", "Asia; Other", "260", "2600", "115.00", "24"],
        ["Overall total", "-", "-", "-", "-", "1050", "10500", "525.00", "103"],
    ]
    with path.open("w", encoding="utf-8", newline="") as handle:
        csv.writer(handle).writerows(rows)


def _write_google_campaign_export(path: Path) -> None:
    rows = [
        ["Campaign", "Labels", "Clicks", "Impr.", "Cost", "Conversions"],
        ["UK - Generic - South Korea - General", "Asia; Other", "180", "1800", "95.00", "25"],
        ["UK - Generic - China - General", "Asia; China", "170", "1700", "130.00", "40"],
    ]
    with path.open("w", encoding="utf-8", newline="") as handle:
        csv.writer(handle).writerows(rows)


def _write_google_time_series_export(path: Path) -> None:
    rows = [["Date", "Impr."], ["Wed, 1 Apr 2026", "54,238"], ["Thu, 2 Apr 2026", "75,209"]]
    with path.open("w", encoding="utf-8", newline="") as handle:
        csv.writer(handle).writerows(rows)


def _write_performance_csv(path: Path) -> Path:
    rows = []
    for month in (4, 5, 6):
        rows.extend(
            [
                {
                    "Date": f"01/{month:02d}/2026",
                    "Campaign Type": "Brand",
                    "Destination": "China",
                    "Impressions": 1000,
                    "Clicks": 100,
                    "Cost": 50,
                    "Sales Leads": 20,
                },
                {
                    "Date": f"01/{month:02d}/2026",
                    "Campaign Type": "Generic",
                    "Destination": "Other",
                    "Impressions": 2000,
                    "Clicks": 200,
                    "Cost": 100,
                    "Sales Leads": 15,
                },
            ]
        )
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)
    return path


def _pptx_text(pptx_path: Path) -> str:
    presentation = Presentation(str(pptx_path))
    chunks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
    return "\n".join(chunks)
