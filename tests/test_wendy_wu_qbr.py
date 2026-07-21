from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation

from src.chart_builder import ChartBuilder
from src.config_loader import ConfigLoader
from src.data_loader import detect_latest_complete_month, detect_latest_complete_quarter, load_csv
from src.metrics import prepare_report_data, validate_report_data
from src.report_pipeline import ReportPipeline
from utils.text_report import TextReportPipeline, generate_text_report


ROOT = Path(__file__).resolve().parent.parent
CONFIG_LOADER = ConfigLoader(
    report_config_path=ROOT / "config" / "report_config.yaml",
    chart_styles_path=ROOT / "config" / "chart_styles.yaml",
    clients_config_path=ROOT / "config" / "clients_config.json",
)
FIXTURES = {
    "wendy_wu": ROOT / "temp_uploads/9dbf730e0cb845faaff0d726d4cefc7b/performance/Wendy Wu UK data.csv",
    "wendy_wu_australia": ROOT / "temp_uploads/517d51c54ffa403599f70e83031af90b/performance/Wendy Wu AU data.csv",
}
EXPECTED_KPI_LABELS = ["Cost", "Sales Leads", "CPL", "CVR", "Clicks", "Impressions", "CTR", "CPC"]
SECTION_DIVIDER = "-" * 40


class WendyWuQbrTests(unittest.TestCase):
    def test_other_rollup_excludes_brand_and_matches_non_brand_remainder(self) -> None:
        for client_id, csv_path in FIXTURES.items():
            with self.subTest(client_id=client_id):
                client_config, quarter, report, df = _prepare_report(client_id, csv_path)
                expected_df = _apply_destination_aliases_for_test(df, client_config)
                current_df = expected_df[(expected_df["year"] == quarter.year) & (expected_df["quarter"] == quarter.quarter)].copy()

                self.assertIn("Other", report["available_destinations"])
                self.assertNotIn("Brand", report["dest_mix"]["Other"]["Campaign Type"].tolist())

                non_brand_remainder = current_df[
                    (~current_df["destination"].isin(client_config["destinations"]))
                    & (current_df["campaign_type"] != "Brand")
                ].copy()
                other_total = report["destinations"]["Other"]["total"]

                self.assertAlmostEqual(other_total["Cost"], float(non_brand_remainder["cost"].sum()))
                self.assertAlmostEqual(other_total["Sales Leads"], float(non_brand_remainder["sales_leads"].sum()))
                self.assertAlmostEqual(other_total["Clicks"], float(non_brand_remainder["clicks"].sum()))

                excluded_total = report["destination_excluded_total"]
                excluded_brand = current_df[
                    (~current_df["destination"].isin(client_config["destinations"]))
                    & (current_df["campaign_type"] == "Brand")
                ].copy()
                self.assertAlmostEqual(excluded_total["Cost"], float(excluded_brand["cost"].sum()))
                self.assertAlmostEqual(excluded_total["Sales Leads"], float(excluded_brand["sales_leads"].sum()))

    def test_other_rollup_absorbs_non_core_destinations_beyond_literal_other(self) -> None:
        df = pd.DataFrame(
            [
                {"Date": "01/01/2026", "Campaign Type": "Brand", "Destination": "Vietnam", "Impressions": 100, "Clicks": 10, "Cost": 20, "Sales Leads": 5},
                {"Date": "02/01/2026", "Campaign Type": "Generic", "Destination": "Vietnam", "Impressions": 120, "Clicks": 12, "Cost": 24, "Sales Leads": 4},
                {"Date": "03/01/2026", "Campaign Type": "Generic", "Destination": "Other", "Impressions": 140, "Clicks": 14, "Cost": 28, "Sales Leads": 7},
                {"Date": "04/02/2026", "Campaign Type": "Generic", "Destination": "China", "Impressions": 160, "Clicks": 16, "Cost": 32, "Sales Leads": 8},
                {"Date": "05/03/2026", "Campaign Type": "Generic", "Destination": "Laos", "Impressions": 180, "Clicks": 18, "Cost": 36, "Sales Leads": 9},
            ]
        )
        with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as tmp:
            df.to_csv(tmp.name, index=False)
            csv_path = Path(tmp.name)

        loaded = load_csv(csv_path)
        quarter = detect_latest_complete_quarter(loaded)
        report = prepare_report_data(
            loaded,
            quarter,
            campaign_order=["Brand", "Generic"],
            destination_order=["China"],
            destination_other_config={
                "enabled": True,
                "label": "Other",
                "mode": "remainder",
                "exclude_campaign_types": ["Brand"],
            },
        )

        other_total = report["destinations"]["Other"]["total"]
        self.assertEqual(report["available_destinations"], ["China", "Other"])
        self.assertAlmostEqual(other_total["Sales Leads"], 20.0)
        self.assertAlmostEqual(other_total["Cost"], 88.0)
        self.assertNotIn("Brand", report["dest_mix"]["Other"]["Campaign Type"].tolist())

    def test_named_destination_totals_and_kpi_payloads_are_stable(self) -> None:
        for client_id, csv_path in FIXTURES.items():
            with self.subTest(client_id=client_id):
                client_config, quarter, report, df = _prepare_report(client_id, csv_path)
                expected_df = _apply_destination_aliases_for_test(df, client_config)
                current_df = expected_df[(expected_df["year"] == quarter.year) & (expected_df["quarter"] == quarter.quarter)].copy()

                overall_raw_cost = float(current_df["cost"].sum())
                campaign_cost = sum(scope["total"]["Cost"] for scope in report["campaigns"].values())
                self.assertAlmostEqual(report["overall"]["total"]["Cost"], overall_raw_cost)
                self.assertAlmostEqual(campaign_cost, overall_raw_cost)

                for destination in client_config["destinations"]:
                    raw_subset = current_df[current_df["destination"] == destination].copy()
                    if raw_subset.empty:
                        self.assertNotIn(destination, report["destinations"])
                        continue
                    destination_scope = report["destinations"][destination]
                    self.assertAlmostEqual(destination_scope["total"]["Cost"], float(raw_subset["cost"].sum()))
                    self.assertAlmostEqual(destination_scope["total"]["Sales Leads"], float(raw_subset["sales_leads"].sum()))

                for scope in [report["overall"], report["campaigns"]["Brand"], report["destinations"]["China"]]:
                    self.assertEqual([item["label"] for item in scope["kpis"]], EXPECTED_KPI_LABELS)
                    self.assertEqual(len(scope["kpis"]), len(EXPECTED_KPI_LABELS))

    def test_kpi_yoy_falls_back_to_na_when_prior_year_baseline_is_missing(self) -> None:
        df = load_csv(_write_single_quarter_fixture())
        quarter = detect_latest_complete_quarter(df)
        report = prepare_report_data(
            df,
            quarter,
            campaign_order=["Brand"],
            destination_order=["China"],
            destination_other_config={"enabled": True, "label": "Other", "mode": "remainder", "exclude_campaign_types": ["Brand"]},
        )
        validate_report_data(report)

        for kpi in report["overall"]["kpis"]:
            self.assertIsNone(kpi["yoy"])
            self.assertEqual(kpi["yoy_label"], "n/a")

    def test_yoy_falls_back_to_na_when_current_derived_kpi_is_undefined(self) -> None:
        rows = [
            {"Date": "01/04/2025", "Campaign Type": "Other", "Destination": "Other", "Impressions": 100, "Clicks": 10, "Cost": 20, "Sales Leads": 2},
            {"Date": "01/05/2025", "Campaign Type": "Other", "Destination": "Other", "Impressions": 100, "Clicks": 10, "Cost": 20, "Sales Leads": 2},
            {"Date": "01/06/2025", "Campaign Type": "Other", "Destination": "Other", "Impressions": 100, "Clicks": 10, "Cost": 20, "Sales Leads": 2},
            {"Date": "01/04/2026", "Campaign Type": "Other", "Destination": "Other", "Impressions": 0, "Clicks": 0, "Cost": 0, "Sales Leads": 1},
            {"Date": "01/05/2026", "Campaign Type": "Other", "Destination": "Other", "Impressions": 0, "Clicks": 0, "Cost": 0, "Sales Leads": 1},
            {"Date": "01/06/2026", "Campaign Type": "Other", "Destination": "Other", "Impressions": 0, "Clicks": 0, "Cost": 0, "Sales Leads": 1},
        ]
        with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as tmp:
            pd.DataFrame(rows).to_csv(tmp.name, index=False)
            csv_path = Path(tmp.name)

        df = load_csv(csv_path)
        quarter = detect_latest_complete_quarter(df)
        report = prepare_report_data(
            df,
            quarter,
            campaign_order=["Brand", "Generic"],
            destination_order=["China"],
            destination_other_config={"enabled": True, "label": "Other", "mode": "remainder", "exclude_campaign_types": ["Brand"]},
        )
        validate_report_data(report)

        other_kpis = {kpi["key"]: kpi for kpi in report["campaigns"]["Other"]["kpis"]}
        self.assertIsNone(report["campaigns"]["Other"]["yoy"]["CTR"])
        self.assertIsNone(report["campaigns"]["Other"]["yoy"]["CPC"])
        self.assertIsNone(report["campaigns"]["Other"]["yoy"]["CVR"])
        self.assertEqual(other_kpis["CTR"]["yoy_label"], "n/a")
        self.assertEqual(other_kpis["CPC"]["yoy_label"], "n/a")
        self.assertEqual(other_kpis["CVR"]["yoy_label"], "n/a")

    def test_text_month_dates_are_parsed_across_quarter_months(self) -> None:
        rows = [
            {"Date": "1 Jan 2026", "Campaign Type": "Brand", "Destination": "China", "Impressions": 100, "Clicks": 10, "Cost": 20, "Sales Leads": 5},
            {"Date": "1 Feb 2026", "Campaign Type": "Brand", "Destination": "China", "Impressions": 120, "Clicks": 12, "Cost": 24, "Sales Leads": 6},
            {"Date": "1 Mar 2026", "Campaign Type": "Brand", "Destination": "China", "Impressions": 140, "Clicks": 14, "Cost": 28, "Sales Leads": 7},
            {"Date": "1 May 2026", "Campaign Type": "Brand", "Destination": "China", "Impressions": 160, "Clicks": 16, "Cost": 32, "Sales Leads": 8},
        ]
        frame = pd.DataFrame(rows)
        with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as tmp:
            frame.to_csv(tmp.name, index=False)
            csv_path = Path(tmp.name)

        df = load_csv(csv_path)
        quarter = detect_latest_complete_quarter(df)

        self.assertEqual(len(df), 4)
        self.assertEqual((quarter.year, quarter.quarter), (2026, 1))

    def test_text_report_includes_kpi_block_and_other_mix_excludes_brand(self) -> None:
        client_id = "wendy_wu"
        client_config, quarter, report, _ = _prepare_report(client_id, FIXTURES[client_id])
        report_text = generate_text_report(
            report,
            campaigns=report["available_campaigns"],
            other_inputs_if_needed={
                "client_name": client_config["name"],
                "report_title": client_config["report_title"],
                "agency_name": client_config["agency"],
                "subtitle": f"{quarter.label} ({quarter.start.strftime('%b')} - {quarter.end.strftime('%b %Y')})",
                "client_config": client_config,
                "config_loader": CONFIG_LOADER,
                "trends_summary": None,
                "auction_summary": None,
                "recommendations": [],
            },
        )

        self.assertIn("Key Metrics + YoY", report_text)
        other_mix_section = _extract_section(report_text, "Other Campaign Mix")
        self.assertNotIn("Brand", other_mix_section)

        other_trend_section = _extract_section(report_text, "Other Monthly Trend")
        self.assertNotIn("- Other generated", other_trend_section)

    def test_chart_builder_outputs_dual_axis_wendy_figures(self) -> None:
        for client_id, csv_path in FIXTURES.items():
            with self.subTest(client_id=client_id):
                client_config, _, report, _ = _prepare_report(client_id, csv_path)
                chart_styles = CONFIG_LOADER.get_chart_styles(client_config)
                with tempfile.TemporaryDirectory() as tmpdir:
                    chart_builder = ChartBuilder(Path(tmpdir), chart_styles=chart_styles)
                    charts = chart_builder.build_scope_trend_charts("other_scope", report["destinations"]["Other"]["monthly"])
                    self.assertTrue(charts["cpl_cvr"].exists())
                    self.assertTrue(charts["cost_leads"].exists())

                    fig, ax1, ax2 = chart_builder.build_cpl_cvr_figure(report["destinations"]["Other"]["monthly"])
                    self.assertEqual(ax1.get_xlabel(), "Month")
                    self.assertEqual(ax1.get_ylabel(), "CPL (£)")
                    self.assertEqual(ax2.get_ylabel(), "CVR (%)")
                    self.assertEqual(len(ax1.lines), 1)
                    self.assertEqual(len(ax2.lines), 1)
                    self.assertEqual(ax2.lines[0].get_color(), client_config["branding"]["chart_palette"]["cvr"])
                    plt.close(fig)

                    fig, cost_ax1, cost_ax2 = chart_builder.build_cost_leads_figure(report["destinations"]["Other"]["monthly"])
                    self.assertEqual(cost_ax1.get_xlabel(), "Month")
                    self.assertEqual(cost_ax1.get_ylabel(), "Cost (£)")
                    self.assertEqual(cost_ax2.get_ylabel(), "Sales Leads")
                    self.assertEqual(len(cost_ax2.lines), 1)
                    self.assertGreaterEqual(len(cost_ax1.containers), 1)
                    plt.close(fig)

    def test_central_asia_mongolia_is_uk_only_named_destination(self) -> None:
        csv_path = _write_central_asia_fixture()
        uk_config = CONFIG_LOADER.get_client_config("wendy_wu")
        australia_config = CONFIG_LOADER.get_client_config("wendy_wu_australia")

        self.assertIn("Central Asia & Mongolia", uk_config["destinations"])
        self.assertNotIn("Central Asia & Mongolia", australia_config["destinations"])

        _, _, uk_report, _ = _prepare_report("wendy_wu", csv_path)
        _, _, australia_report, _ = _prepare_report("wendy_wu_australia", csv_path)

        self.assertIn("Central Asia & Mongolia", uk_report["available_destinations"])
        self.assertIn("Other", uk_report["available_destinations"])
        self.assertAlmostEqual(uk_report["destinations"]["Central Asia & Mongolia"]["total"]["Cost"], 600.0)
        self.assertAlmostEqual(uk_report["destinations"]["Other"]["total"]["Cost"], 300.0)

        self.assertNotIn("Central Asia & Mongolia", australia_report["available_destinations"])
        self.assertAlmostEqual(australia_report["destinations"]["Other"]["total"]["Cost"], 900.0)

    def test_destination_campaign_mix_includes_inline_yoy_for_all_visible_metrics(self) -> None:
        csv_path = _write_destination_mix_yoy_fixture()
        client_config, quarter, report, _ = _prepare_report("wendy_wu_australia", csv_path)

        report_text = generate_text_report(
            report,
            campaigns=report["available_campaigns"],
            other_inputs_if_needed={
                "client_name": client_config["name"],
                "report_title": client_config["report_title"],
                "agency_name": client_config["agency"],
                "subtitle": f"{quarter.label} ({quarter.start.strftime('%b')} - {quarter.end.strftime('%b %Y')})",
                "client_config": client_config,
                "config_loader": CONFIG_LOADER,
                "trends_summary": None,
                "auction_summary": None,
                "recommendations": [],
            },
        )

        mix_section = _extract_section(report_text, "China Campaign Mix")
        self.assertIn("£600.00 (+100%)", mix_section)
        self.assertIn("60 (+100%)", mix_section)
        self.assertIn("66.67% (+33%)", mix_section)
        self.assertIn("66.67% (+100%)", mix_section)
        self.assertIn("£10.00 (+0%)", mix_section)

    def test_uk_pptx_adds_central_asia_mongolia_destination_slides(self) -> None:
        csv_path = _write_central_asia_fixture()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "wendy_wu_uk.pptx"
            pipeline = ReportPipeline(project_root=ROOT)
            pipeline.charts_root = Path(tmpdir) / "charts"
            generated_path = pipeline.run(
                input_csv=csv_path,
                output_pptx=output_path,
                client_id="wendy_wu",
            )

            titles = _pptx_text(generated_path)
            self.assertIn("Central Asia & Mongolia Summary + YoY", titles)
            self.assertIn("Central Asia & Mongolia Monthly Trend", titles)
            self.assertIn("Central Asia & Mongolia Campaign Mix", titles)

    def test_uk_qbr_text_report_uses_ytd_trend_pair_uploads(self) -> None:
        csv_path = _write_monthly_fixture()
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            current_trends = root / "current_trends"
            previous_trends = root / "previous_trends"
            current_trends.mkdir()
            previous_trends.mkdir()
            _write_trend_csv(current_trends / "brand_current.csv", "wendy wu tours", 2026, [10, 20, 30, 40, 50, 60])
            _write_trend_csv(previous_trends / "brand_previous.csv", "wendy wu tours", 2025, [5, 10, 15, 20, 25, 30])

            report_txt = root / "report.txt"
            TextReportPipeline(project_root=ROOT).run(
                input_csv=csv_path,
                output_txt=report_txt,
                client_id="wendy_wu",
                trends_ytd_current_dir=current_trends,
                trends_ytd_previous_dir=previous_trends,
            )

            report_text = report_txt.read_text(encoding="utf-8")

        self.assertIn("Brand YTD demand has increased", report_text)
        self.assertIn("Jan", report_text)
        self.assertIn("Jun", report_text)
        self.assertIn("60.0", report_text)
        self.assertIn("30.0", report_text)

    def test_australia_qbr_text_report_uses_ytd_trend_pair_uploads(self) -> None:
        csv_path = _write_monthly_fixture()
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            current_trends = root / "current_trends"
            previous_trends = root / "previous_trends"
            current_trends.mkdir()
            previous_trends.mkdir()
            _write_trend_csv(current_trends / "brand_current.csv", "wendy wu tours", 2026, [12, 18, 24, 30, 36, 42])
            _write_trend_csv(previous_trends / "brand_previous.csv", "wendy wu tours", 2025, [6, 9, 12, 15, 18, 21])

            report_txt = root / "report.txt"
            TextReportPipeline(project_root=ROOT).run(
                input_csv=csv_path,
                output_txt=report_txt,
                client_id="wendy_wu_australia",
                trends_ytd_current_dir=current_trends,
                trends_ytd_previous_dir=previous_trends,
            )

            report_text = report_txt.read_text(encoding="utf-8")

        self.assertIn("Brand YTD demand has increased", report_text)
        self.assertIn("Jan", report_text)
        self.assertIn("Jun", report_text)
        self.assertIn("42.0", report_text)
        self.assertIn("21.0", report_text)

    def test_monthly_report_uses_latest_full_month_cards_and_ytd_tables(self) -> None:
        csv_path = _write_monthly_fixture()
        client_config = CONFIG_LOADER.get_client_config("wendy_wu")
        df = load_csv(csv_path)
        month = detect_latest_complete_month(df, today=pd.Timestamp("2026-07-07"))
        self.assertEqual((month.year, month.month), (2026, 6))

        report = prepare_report_data(
            df,
            month,
            campaign_order=CONFIG_LOADER.get_campaign_types(client_config),
            destination_order=CONFIG_LOADER.get_destinations(client_config),
            destination_aliases=client_config.get("destination_aliases"),
            destination_other_config=client_config.get("destination_other"),
            report_mode="monthly",
        )
        validate_report_data(report)

        june_df = df[(df["year"] == 2026) & (df["month"] == 6)].copy()
        ytd_df = df[(df["year"] == 2026) & (df["month"] <= 6)].copy()
        self.assertAlmostEqual(report["overall"]["total"]["Cost"], float(june_df["cost"].sum()))
        total_row = report["overall"]["monthly"][report["overall"]["monthly"]["Month"] == "Total"].iloc[0]
        self.assertAlmostEqual(total_row["Cost"], float(ytd_df["cost"].sum()))
        self.assertEqual(report["overall"]["monthly"]["Month"].tolist(), ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Total"])

        kpis = {kpi["key"]: kpi for kpi in report["overall"]["kpis"]}
        self.assertIn("mom_label", kpis["Cost"])
        self.assertIn("yoy_label", kpis["Cost"])

    def test_monthly_pptx_uses_performance_only_sections(self) -> None:
        csv_path = _write_monthly_fixture()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "wendy_wu_monthly.pptx"
            pipeline = ReportPipeline(project_root=ROOT)
            pipeline.charts_root = Path(tmpdir) / "charts"
            generated_path = pipeline.run(
                input_csv=csv_path,
                output_pptx=output_path,
                client_id="wendy_wu",
                report_mode="monthly",
            )

            titles = _pptx_text(generated_path)
            self.assertIn("Overall Month Summary", titles)
            self.assertIn("Campaign Type YTD Mix", titles)
            self.assertIn("China Month Summary + YoY", titles)
            self.assertIn("China YTD Trend", titles)
            self.assertNotIn("Google Trends", titles)
            self.assertNotIn("Auction Insights", titles)
            self.assertNotIn("Recommendations", titles)


def _prepare_report(client_id: str, csv_path: Path):
    client_config = CONFIG_LOADER.get_client_config(client_id)
    df = load_csv(csv_path)
    quarter = detect_latest_complete_quarter(df)
    report = prepare_report_data(
        df,
        quarter,
        campaign_order=CONFIG_LOADER.get_campaign_types(client_config),
        destination_order=CONFIG_LOADER.get_destinations(client_config),
        destination_aliases=client_config.get("destination_aliases"),
        destination_other_config=client_config.get("destination_other"),
    )
    validate_report_data(report)
    return client_config, quarter, report, df


def _apply_destination_aliases_for_test(df: pd.DataFrame, client_config: dict) -> pd.DataFrame:
    aliases = client_config.get("destination_aliases") or {}
    if not aliases:
        return df.copy()

    lookup = {}
    for canonical, alias_values in aliases.items():
        values = [canonical, *(alias_values or [])]
        for value in values:
            lookup[_normalize_destination_for_test(value)] = canonical

    expected = df.copy()
    expected["destination"] = expected["destination"].map(
        lambda value: lookup.get(_normalize_destination_for_test(value), str(value).strip())
    )
    return expected


def _normalize_destination_for_test(value: object) -> str:
    return "".join(ch for ch in str(value).strip().lower() if ch.isalnum())


def _extract_section(report_text: str, section_title: str) -> str:
    start = report_text.find(section_title)
    if start == -1:
        return ""
    next_divider = report_text.find(SECTION_DIVIDER, start + len(section_title))
    return report_text[start:] if next_divider == -1 else report_text[start:next_divider]


def _write_single_quarter_fixture() -> Path:
    rows = [
        {"Date": "01/01/2026", "Campaign Type": "Brand", "Destination": "China", "Impressions": 120, "Clicks": 12, "Cost": 24, "Sales Leads": 6},
        {"Date": "01/02/2026", "Campaign Type": "Brand", "Destination": "China", "Impressions": 100, "Clicks": 10, "Cost": 20, "Sales Leads": 5},
        {"Date": "01/03/2026", "Campaign Type": "Brand", "Destination": "China", "Impressions": 80, "Clicks": 8, "Cost": 16, "Sales Leads": 4},
    ]
    frame = pd.DataFrame(rows)
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as tmp:
        frame.to_csv(tmp.name, index=False)
    return Path(tmp.name)


def _write_central_asia_fixture() -> Path:
    rows = []
    for year, factor in [(2025, 0.5), (2026, 1.0)]:
        for month, central_cost, other_cost, central_destination in [
            (1, 100, 50, "Central Asia"),
            (2, 200, 100, "Mongolia"),
            (3, 300, 150, "Central Asia & Mongolia"),
        ]:
            rows.extend(
                [
                    {
                        "Date": f"01/{month:02d}/{year}",
                        "Campaign Type": "Generic",
                        "Destination": "China",
                        "Impressions": 1000,
                        "Clicks": 100,
                        "Cost": 25 * factor,
                        "Sales Leads": 5 * factor,
                    },
                    {
                        "Date": f"01/{month:02d}/{year}",
                        "Campaign Type": "Generic",
                        "Destination": central_destination,
                        "Impressions": 2000,
                        "Clicks": 200,
                        "Cost": central_cost * factor,
                        "Sales Leads": 20 * factor,
                    },
                    {
                        "Date": f"01/{month:02d}/{year}",
                        "Campaign Type": "Generic",
                        "Destination": "Other",
                        "Impressions": 1500,
                        "Clicks": 150,
                        "Cost": other_cost * factor,
                        "Sales Leads": 10 * factor,
                    },
                ]
            )
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as tmp:
        pd.DataFrame(rows).to_csv(tmp.name, index=False)
        return Path(tmp.name)


def _write_destination_mix_yoy_fixture() -> Path:
    rows = []
    for month in (1, 2, 3):
        rows.extend(
            [
                {"Date": f"01/{month:02d}/2025", "Campaign Type": "Generic", "Destination": "China", "Impressions": 1000, "Clicks": 100, "Cost": 100, "Sales Leads": 10},
                {"Date": f"01/{month:02d}/2025", "Campaign Type": "Demand Gen", "Destination": "China", "Impressions": 1000, "Clicks": 100, "Cost": 100, "Sales Leads": 20},
                {"Date": f"01/{month:02d}/2026", "Campaign Type": "Generic", "Destination": "China", "Impressions": 1000, "Clicks": 100, "Cost": 200, "Sales Leads": 20},
                {"Date": f"01/{month:02d}/2026", "Campaign Type": "Demand Gen", "Destination": "China", "Impressions": 1000, "Clicks": 100, "Cost": 100, "Sales Leads": 10},
            ]
        )
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as tmp:
        pd.DataFrame(rows).to_csv(tmp.name, index=False)
        return Path(tmp.name)


def _write_monthly_fixture() -> Path:
    rows = []
    for year in (2025, 2026):
        for month in range(1, 8):
            rows.extend(
                [
                    {
                        "Date": f"01/{month:02d}/{year}",
                        "Campaign Type": "Brand",
                        "Destination": "China",
                        "Impressions": 1000 + month,
                        "Clicks": 100 + month,
                        "Cost": 100 + month * (2 if year == 2026 else 1),
                        "Sales Leads": 10 + month,
                    },
                    {
                        "Date": f"01/{month:02d}/{year}",
                        "Campaign Type": "Generic",
                        "Destination": "Central Asia & Mongolia",
                        "Impressions": 800 + month,
                        "Clicks": 80 + month,
                        "Cost": 80 + month * (3 if year == 2026 else 1),
                        "Sales Leads": 8 + month,
                    },
                    {
                        "Date": f"01/{month:02d}/{year}",
                        "Campaign Type": "Generic",
                        "Destination": "Other",
                        "Impressions": 600 + month,
                        "Clicks": 60 + month,
                        "Cost": 60 + month,
                        "Sales Leads": 6 + month,
                    },
                ]
            )
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as tmp:
        pd.DataFrame(rows).to_csv(tmp.name, index=False)
        return Path(tmp.name)


def _write_trend_csv(path: Path, term: str, year: int, values: list[int]) -> None:
    rows = [{"Date": f"{year}-{month:02d}-01", term: value} for month, value in enumerate(values, start=1)]
    pd.DataFrame(rows).to_csv(path, index=False)


def _pptx_text(pptx_path: Path) -> str:
    presentation = Presentation(str(pptx_path))
    chunks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
    return "\n".join(chunks)


if __name__ == "__main__":
    unittest.main()
